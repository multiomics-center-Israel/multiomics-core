#!/usr/bin/env python3
"""Add roadmap/flowchart slides to existing summary PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Colour palette (same as main deck)
# ---------------------------------------------------------------------------
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD   = RGBColor(0x1B, 0x2A, 0x4A)
BG_CARD2  = RGBColor(0x16, 0x23, 0x3D)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)  # sky-400
ACCENT2   = RGBColor(0x81, 0x8C, 0xF8)  # indigo-400
GREEN     = RGBColor(0x4A, 0xDE, 0x80)  # green-400
SUBTLE    = RGBColor(0x94, 0xA3, 0xB8)  # slate-400
ORANGE    = RGBColor(0xFB, 0x92, 0x3C)  # orange-400
PINK      = RGBColor(0xF4, 0x72, 0xB6)  # pink-400
YELLOW    = RGBColor(0xFA, 0xCC, 0x15)  # yellow-400
DIM       = RGBColor(0x47, 0x55, 0x69)  # slate-600
DONE_BG   = RGBColor(0x0D, 0x2B, 0x1A)  # dark green tint
PROG_BG   = RGBColor(0x2B, 0x1D, 0x0D)  # dark amber tint
PLAN_BG   = RGBColor(0x1A, 0x1A, 0x2E)  # dark indigo tint

BASE = os.path.dirname(os.path.abspath(__file__))
prs = Presentation(f"{BASE}/multiomics_core_summary.pptx")

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, size=16, color=WHITE, bold=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    return p

def card_rect(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def circle(slide, left, top, size, fill_color, text="", text_size=12, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    return shape

def arrow_right(slide, left, top, width, color=SUBTLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def arrow_down(slide, left, top, height, color=SUBTLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.25), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def footer(slide, text="Multi-Omics Core  |  Dor Branch Summary  |  2026"):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             text, size=9, color=SUBTLE, alignment=PP_ALIGN.LEFT)

def status_badge(slide, left, top, text, color):
    """Small rounded badge with status text."""
    w = Inches(0.9)
    h = Inches(0.28)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


# =========================================================================
# SLIDE A — DEVELOPMENT JOURNEY (horizontal timeline)
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
         "Development Journey", size=32, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.8), Inches(5), ACCENT)

# --- Horizontal timeline line ---
timeline_y = Inches(2.05)
shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(0.8), timeline_y, Inches(11.7), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = DIM
shape.line.fill.background()

# Timeline phases
phases = [
    {
        "label": "Phase 1",
        "title": "Foundation",
        "color": ACCENT,
        "items": [
            "targets orchestration",
            "Config system (YAML)",
            "Core plot library",
            "Excel export engine",
        ],
        "status": "DONE",
    },
    {
        "label": "Phase 2",
        "title": "Proteomics",
        "color": ACCENT,
        "items": [
            "MaxQuant/DIA-NN input",
            "QRILC imputation + VSN",
            "Limma DE + fdrtool",
            "PPI networks (STRING)",
        ],
        "status": "DONE",
    },
    {
        "label": "Phase 3",
        "title": "RNA-seq",
        "color": GREEN,
        "items": [
            "DESeq2 factory",
            "Batch correction",
            "Deconvolution (xCell2)",
            "Binary clustering",
        ],
        "status": "DONE",
    },
    {
        "label": "Phase 4",
        "title": "Metabolomics",
        "color": ORANGE,
        "items": [
            "glog10 + PQN norm",
            "HMDB annotation",
            "RF + PLS-DA selection",
            "4 enrichment methods",
        ],
        "status": "DONE",
    },
    {
        "label": "Phase 5",
        "title": "Reporting &\nAutomation",
        "color": PINK,
        "items": [
            "HTML reports (Rmd)",
            "Auto PPTX + AI",
            "Pipeline summary HTML",
            "Wizard + CLI (run.R)",
        ],
        "status": "DONE",
    },
]

node_w = Inches(2.1)
node_h = Inches(4.2)
start_x = Inches(0.5)
gap = Inches(0.3)

for i, phase in enumerate(phases):
    x = start_x + i * (node_w + gap)
    dot_size = Inches(0.35)

    # Timeline dot
    circle(sl, x + node_w / 2 - dot_size / 2,
           timeline_y - dot_size / 2 + Pt(2),
           dot_size, phase["color"])

    # Phase label above timeline
    add_text(sl, x, Inches(1.15), node_w, Inches(0.3),
             phase["label"], size=10, color=SUBTLE,
             alignment=PP_ALIGN.CENTER, bold=True)

    # Card below timeline
    card_y = Inches(2.55)
    card_rect(sl, x, card_y, node_w, node_h)
    accent_line(sl, x + Inches(0.1), card_y + Inches(0.08), node_w - Inches(0.2), phase["color"])

    add_text(sl, x + Inches(0.1), card_y + Inches(0.2), node_w - Inches(0.2), Inches(0.55),
             phase["title"], size=15, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

    # Items
    tf = add_text(sl, x + Inches(0.15), card_y + Inches(0.85),
                  node_w - Inches(0.3), Inches(2.5),
                  "", size=11, color=SUBTLE)
    for item in phase["items"]:
        add_para(tf, f"  • {item}", size=11, color=SUBTLE, space_before=Pt(5))

    # Status badge
    status_badge(sl, x + node_w / 2 - Inches(0.45),
                 card_y + node_h - Inches(0.45),
                 phase["status"], GREEN)

footer(sl)


# =========================================================================
# SLIDE B — WHERE WE ARE NOW (status dashboard)
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
         "Where We Are Now", size=32, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.8), Inches(5), GREEN)

# === LEFT COLUMN: Pipeline status ===
domains = [
    ("Proteomics Pipeline", ACCENT, "100%", [
        ("Data Import (DIANN/MaxQuant/FragPipe)", GREEN),
        ("Imputation (QRILC/MinVal/KNN)", GREEN),
        ("Normalization (VSN/Median)", GREEN),
        ("DE (Limma + t-test + ANOVA)", GREEN),
        ("Pathway Enrichment (GSEA+ORA)", GREEN),
        ("PPI Networks (STRING)", GREEN),
        ("HTML Report + PPTX", GREEN),
        ("Shiny Export v2.0", GREEN),
    ]),
    ("RNA-seq Pipeline", GREEN, "100%", [
        ("Data Import (counts + tximport)", GREEN),
        ("Batch Correction (ComBat/SVA/RUV)", GREEN),
        ("DE (DESeq2 factory)", GREEN),
        ("Deconvolution (xCell2)", GREEN),
        ("Pathway Enrichment (fGSEA+ORA)", GREEN),
        ("Clustering (Hierarchical/Binary)", GREEN),
        ("HTML Report + PPTX", GREEN),
        ("Shiny Export v2.0", GREEN),
    ]),
    ("Metabolomics Pipeline", ORANGE, "100%", [
        ("Data Import + HMDB Annotation", GREEN),
        ("Normalization (PQN/glog10)", GREEN),
        ("DE (t-test, raw-p support)", GREEN),
        ("Feature Selection (RF + PLS-DA)", GREEN),
        ("Enrichment (QEA/ssGSEA/ORA/GSEA)", GREEN),
        ("HTML Report + PPTX", GREEN),
        ("Shiny Export v2.0", GREEN),
    ]),
]

for i, (title, title_col, pct, features) in enumerate(domains):
    x = Inches(0.4 + i * 4.2)
    card_rect(sl, x, Inches(1.2), Inches(3.9), Inches(5.6))
    accent_line(sl, x + Inches(0.1), Inches(1.3), Inches(3.7), title_col)

    # Title + percentage
    add_text(sl, x + Inches(0.15), Inches(1.45), Inches(2.8), Inches(0.4),
             title, size=15, color=WHITE, bold=True)
    add_text(sl, x + Inches(2.8), Inches(1.45), Inches(0.9), Inches(0.4),
             pct, size=15, color=GREEN, bold=True, alignment=PP_ALIGN.RIGHT)

    # Progress bar background
    bar_y = Inches(1.92)
    bar_shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x + Inches(0.15), bar_y,
                                     Inches(3.6), Inches(0.18))
    bar_shape.fill.solid()
    bar_shape.fill.fore_color.rgb = DIM
    bar_shape.line.fill.background()

    # Progress bar fill
    fill_shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(0.15), bar_y,
                                      Inches(3.6), Inches(0.18))
    fill_shape.fill.solid()
    fill_shape.fill.fore_color.rgb = GREEN
    fill_shape.line.fill.background()

    # Feature checklist
    tf = add_text(sl, x + Inches(0.15), Inches(2.25),
                  Inches(3.6), Inches(4.2), "", size=11, color=SUBTLE)
    for feat, col in features:
        check = "✓" if col == GREEN else "○"
        add_para(tf, f"  {check}  {feat}", size=11, color=col, space_before=Pt(5))

footer(sl)


# =========================================================================
# SLIDE C — ROADMAP / WHAT'S NEXT (flowchart style)
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.25), Inches(12), Inches(0.6),
         "Roadmap — What's Next", size=32, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.8), Inches(5), ACCENT2)

# Legend
add_text(sl, Inches(8.5), Inches(0.3), Inches(4.5), Inches(0.3),
         "", size=1, color=BG_DARK)
for j, (lbl, col) in enumerate([("Done", GREEN), ("Near-term", ORANGE), ("Future", ACCENT2)]):
    bx = Inches(8.5 + j * 1.6)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, bx, Inches(0.42), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    add_text(sl, bx + Inches(0.25), Inches(0.37), Inches(1.2), Inches(0.3),
             lbl, size=11, color=SUBTLE, bold=True)

# ── ROW 1: DONE ──
row1_y = Inches(1.2)
add_text(sl, Inches(0.4), row1_y, Inches(1.2), Inches(0.3),
         "COMPLETED", size=11, color=GREEN, bold=True)
done_items = [
    ("Proteomics\nPipeline", ACCENT),
    ("RNA-seq\nPipeline", GREEN),
    ("Metabolomics\nPipeline", ORANGE),
    ("HTML\nReports", PINK),
    ("Auto\nPPTX", ACCENT2),
    ("Wizard\n+ CLI", ACCENT),
]

box_w = Inches(1.65)
box_h = Inches(1.15)
row1_card_y = Inches(1.55)

for i, (label, col) in enumerate(done_items):
    x = Inches(0.4 + i * 2.1)
    card_rect(sl, x, row1_card_y, box_w, box_h, DONE_BG)
    accent_line(sl, x + Inches(0.05), row1_card_y + Inches(0.05), box_w - Inches(0.1), col)
    add_text(sl, x + Inches(0.05), row1_card_y + Inches(0.15), box_w - Inches(0.1), Inches(0.7),
             label, size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # check mark
    add_text(sl, x + box_w - Inches(0.35), row1_card_y + Inches(0.7), Inches(0.3), Inches(0.3),
             "✓", size=14, color=GREEN, bold=True, alignment=PP_ALIGN.RIGHT)

# ── Big arrow down ──
arrow_down(sl, Inches(6.3), Inches(2.85), Inches(0.4), DIM)

# ── ROW 2: NEAR-TERM (next steps) ──
row2_y = Inches(3.4)
add_text(sl, Inches(0.4), row2_y, Inches(1.5), Inches(0.3),
         "NEAR-TERM", size=11, color=ORANGE, bold=True)

near_items = [
    ("Shiny\nWeb App", "Deploy interactive\nexplorer from v2.0\nShiny contract", ORANGE),
    ("TSS / LOESS\nNormalization", "Drift correction,\nCLR transform for\nmetabolomics", ORANGE),
    ("CI/CD\nGitHub Actions", "Automated testing\non push/PR,\nDocker build", ORANGE),
    ("Lipidomics\nPipeline", "Lipid-specific\npreprocessing,\nDB integration", ORANGE),
]

row2_card_y = Inches(3.75)
near_w = Inches(2.8)
near_h = Inches(1.65)

for i, (title, desc, col) in enumerate(near_items):
    x = Inches(0.4 + i * 3.15)
    card_rect(sl, x, row2_card_y, near_w, near_h, PROG_BG)
    accent_line(sl, x + Inches(0.08), row2_card_y + Inches(0.06), near_w - Inches(0.16), col)
    add_text(sl, x + Inches(0.1), row2_card_y + Inches(0.15), near_w - Inches(0.2), Inches(0.55),
             title, size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.1), row2_card_y + Inches(0.75), near_w - Inches(0.2), Inches(0.75),
             desc, size=10, color=SUBTLE, alignment=PP_ALIGN.CENTER)

# ── Big arrow down ──
arrow_down(sl, Inches(6.3), Inches(5.55), Inches(0.4), DIM)

# ── ROW 3: FUTURE ──
row3_y = Inches(6.1)
add_text(sl, Inches(0.4), row3_y, Inches(1.2), Inches(0.3),
         "FUTURE", size=11, color=ACCENT2, bold=True)

future_items = [
    ("Multi-Omics Integration", "Cross-domain correlation, joint pathway analysis, data fusion", ACCENT2),
    ("Advanced ML / AI", "Deep learning feature extraction, automated interpretation", ACCENT2),
    ("Clinical Reporting", "Regulatory-ready reports, biomarker panels, clinical decision support", ACCENT2),
]

row3_card_y = Inches(6.35)
fut_w = Inches(3.95)
fut_h = Inches(0.85)

for i, (title, desc, col) in enumerate(future_items):
    x = Inches(0.4 + i * 4.2)
    card_rect(sl, x, row3_card_y, fut_w, fut_h, PLAN_BG)
    accent_line(sl, x + Inches(0.08), row3_card_y + Inches(0.05), fut_w - Inches(0.16), col)
    add_text(sl, x + Inches(0.1), row3_card_y + Inches(0.12), fut_w - Inches(0.2), Inches(0.3),
             title, size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.1), row3_card_y + Inches(0.45), fut_w - Inches(0.2), Inches(0.35),
             desc, size=9, color=SUBTLE, alignment=PP_ALIGN.CENTER)

footer(sl)


# =========================================================================
# SAVE
# =========================================================================
out_path = f"{BASE}/multiomics_core_summary.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
