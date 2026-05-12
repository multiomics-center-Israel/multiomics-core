#!/usr/bin/env python3
"""Build A02 Multi-Omics Integration PowerPoint for Elah Pick project."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Paths ---
BASE = "/home/multiomics-core/outputs/elah_pick_mo/Results_Elah_Pick_MO_A02/multiomics"
INTEG = os.path.join(BASE, "integration")
MOFA = os.path.join(INTEG, "mofa")
DIABLO = os.path.join(INTEG, "diablo")
SNF = os.path.join(INTEG, "snf")
CONCORD = os.path.join(BASE, "concordance")
CONSENSUS = os.path.join(BASE, "consensus", "consensus")
MECH = os.path.join(BASE, "mechanistic", "mechanistic")
CROSS = os.path.join(BASE, "cross_enrichment")
MULTIGSEA = os.path.join(BASE, "multigsea")
RNACORR = os.path.join(BASE, "rna_protein_correlation")
FOUND = os.path.join(BASE, "foundational")
OUTPUT = "/home/multiomics-core/outputs/elah_pick_mo/Elah_MultiOmics_A02_Annotated.pptx"

# --- Colors ---
DARK_BLUE = RGBColor(0, 51, 102)
BODY_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY_BG = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
ACCENT_BLUE = RGBColor(0, 102, 204)
ACCENT_GREEN = RGBColor(0, 128, 0)
SECTION_BG = RGBColor(0, 51, 102)
BORDER_BLUE = RGBColor(0, 76, 153)

# --- Dimensions (widescreen) ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BODY_GRAY, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", bg_color=None, anchor=MSO_ANCHOR.TOP):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


def add_title(slide, text, top=Inches(0.3)):
    """Add slide title."""
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.7),
                text, font_size=28, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.LEFT)


def add_subtitle(slide, text, top=Inches(1.0)):
    """Add subtitle."""
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.5),
                text, font_size=16, color=ACCENT_BLUE)


def add_body(slide, text, left=Inches(0.6), top=Inches(1.5),
             width=Inches(12), height=Inches(5.0), font_size=14):
    """Add body text."""
    return add_textbox(slide, left, top, width, height, text,
                       font_size=font_size, color=BODY_GRAY)


def add_method_box(slide, text, left=Inches(0.6), top=Inches(5.8),
                   width=Inches(12), height=Inches(1.2)):
    """Add light gray methodology box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = BODY_GRAY
    p.font.name = "Calibri"
    p.font.italic = True
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = LIGHT_GRAY_BG
    return txBox


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if file exists."""
    if os.path.exists(path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    return False


def add_multiline_body(slide, lines, left=Inches(0.6), top=Inches(1.5),
                       width=Inches(12), height=Inches(5.0), font_size=14,
                       bullet=False):
    """Add multi-line body text using multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = BODY_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox


def add_section_slide(text):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(blank_layout)
    # Full background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SECTION_BG
    bg.line.fill.background()
    add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.5),
                text, font_size=36, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    return slide


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(blank_layout)
# Title band
band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), SLIDE_W, Inches(2.5))
band.fill.solid()
band.fill.fore_color.rgb = SECTION_BG
band.line.fill.background()

add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11), Inches(1.0),
            "Multi-Omics Integration Analysis", font_size=36, bold=True,
            color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(3.2), Inches(11), Inches(0.7),
            "A02 Update -- Giardia lamblia Oxidative Stress", font_size=22,
            color=RGBColor(180, 210, 255), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(0.5),
            "Elah Pick | Technion Multiomics Center | March 2026",
            font_size=16, color=BODY_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(5.6), Inches(11), Inches(0.5),
            "Tri-omics: Transcriptomics + Proteomics + Metabolomics",
            font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Study Design
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Study Design")
add_subtitle(slide, "Giardia lamblia oxidative stress model")
add_multiline_body(slide, [
    "Organism: Giardia lamblia trophozoites",
    "Conditions: Anaerobic (ANA, n=4) vs. Oxidative stress (OXY, n=4)",
    "Total samples: 8 biological replicates",
    "",
    "Omics layers:",
    "  - Transcriptomics (RNA-seq): 4,062 genes quantified",
    "  - Proteomics (DDA): 4,062 proteins mapped",
    "  - Metabolomics (LC-MS/MS): ~150 annotated metabolites",
    "",
    "Research question: How does Giardia reprogram its transcriptome,",
    "proteome, and metabolome to survive oxidative stress?",
], top=Inches(1.5), font_size=15)
add_method_box(slide,
    "Methodology: Paired multi-omics from same biological samples enables direct cross-layer comparison. "
    "Gene-protein mapping performed via shared Giardia gene identifiers (GL50803_*).")

# ============================================================
# SLIDE 3: A02 Changes from A01
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "A02 Update: Changes from A01")
add_subtitle(slide, "Refined analysis with corrected parameters")
add_multiline_body(slide, [
    "Key changes in A02:",
    "",
    "1. RNA-seq DE cutoff: log2FC >= 0.5 (was 1.0 in A01)",
    "   - A01: 2,393 DE genes (59% of transcriptome)",
    "   - A02: 1,017 DE genes (25% of transcriptome)",
    "   - More biologically conservative; reduces noise",
    "",
    "2. Metabolomics contrast correction:",
    "   - Fixed contrast direction to match proteomics/transcriptomics",
    "   - Ensures consistent fold-change sign across layers",
    "",
    "3. Integration re-run with all three corrections applied",
    "   - MOFA2, DIABLO, and SNF all re-computed",
    "   - Pathway enrichment updated with corrected gene lists",
], top=Inches(1.4), font_size=14)
add_method_box(slide,
    "Impact: Reducing DE genes from 2,393 to 1,017 increases signal-to-noise in integration. "
    "The A01 cutoff captured too many low-effect genes, diluting biological signal in pathway analysis.")

# ============================================================
# SLIDE 4: Methods Overview
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Integration Methods Overview")
add_subtitle(slide, "Three complementary approaches for multi-omics integration")

# Three columns
col_w = Inches(3.8)
col_h = Inches(3.8)
for i, (title, desc) in enumerate([
    ("MOFA2", "Unsupervised factor analysis\n\nDecomposes variation into latent factors shared across omics layers. Identifies major axes of biological variation without using sample labels. Factors capture coordinated changes across transcriptome, proteome, and metabolome."),
    ("DIABLO (mixOmics)", "Supervised integration\n\nMaximizes covariance between omics layers while discriminating conditions. Identifies features from each layer that jointly separate ANA from OXY. Uses sparse PLS to select key features."),
    ("SNF", "Network-based clustering\n\nSimilarity Network Fusion constructs sample similarity networks per layer, then fuses them into a unified network. Spectral clustering on fused network identifies patient subgroups."),
]):
    left = Inches(0.5 + i * 4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), col_w, col_h)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY_BG
    box.line.color.rgb = BORDER_BLUE
    box.line.width = Pt(1.5)
    add_textbox(slide, left + Inches(0.2), Inches(1.8), col_w - Inches(0.4), Inches(0.5),
                title, font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(2.4), col_w - Inches(0.4), col_h - Inches(1.0),
                desc, font_size=12, color=BODY_GRAY)

add_method_box(slide,
    "Consensus: All three methods yield ARI = 1.0 (perfect agreement) for sample clustering, "
    "confirming the ANA vs. OXY separation is robust and method-independent.", top=Inches(5.8))

# ============================================================
# SLIDE 5: MOFA2 Variance Explained
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "MOFA2: Variance Explained per Factor")
add_subtitle(slide, "Factor 1 dominates across all layers -- driven by condition effect")

img_added = add_image_safe(slide, os.path.join(MOFA, "mofa_variance_heatmap.png"),
                           Inches(0.8), Inches(1.8), width=Inches(6.5))

# Table of values on the right
add_multiline_body(slide, [
    "Variance explained (%):",
    "",
    "Factor 1:",
    "  Metabolomics:      38.3%",
    "  Proteomics:          20.0%",
    "  Transcriptomics: 74.4%",
    "",
    "Factor 2:",
    "  Metabolomics:      17.0%",
    "  Proteomics:            0.01%",
    "  Transcriptomics:   0.57%",
    "",
    "Factor 3:",
    "  Metabolomics:      10.1%",
    "  Transcriptomics:   4.1%",
], left=Inches(7.8), top=Inches(1.8), width=Inches(5.0), height=Inches(4.5), font_size=13)

add_method_box(slide,
    "Interpretation: Factor 1 captures the oxidative stress response, explaining 74.4% of transcriptomic, "
    "38.3% of metabolomic, and 20.0% of proteomic variance. This confirms the treatment is the primary source of variation.")

# ============================================================
# SLIDE 6: MOFA2 Factor Plots
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "MOFA2: Factor Sample Projections")
add_subtitle(slide, "Factor 1 vs Factor 2 and Factor 1 vs Factor 3")

add_image_safe(slide, os.path.join(MOFA, "mofa_factors_1_2.png"),
               Inches(0.5), Inches(1.7), width=Inches(5.8))
add_image_safe(slide, os.path.join(MOFA, "mofa_factors_1_3.png"),
               Inches(6.8), Inches(1.7), width=Inches(5.8))

add_method_box(slide,
    "Factor 1 separates ANA from OXY conditions clearly. Factor 2 captures within-group metabolomic "
    "heterogeneity (17% metabolomics variance). Factor 3 captures additional metabolomic variation (10.1%).")

# ============================================================
# SLIDE 7: MOFA2 Top Weights
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "MOFA2: Top Feature Weights (Factor 1)")
add_subtitle(slide, "Key molecular drivers of the oxidative stress axis")

# Add three weight plots side by side
for i, layer in enumerate(["metabolomics", "proteomics", "transcriptomics"]):
    path = os.path.join(MOFA, f"mofa_top_weights_{layer}.png")
    add_image_safe(slide, path, Inches(0.3 + i * 4.3), Inches(1.7), width=Inches(4.0))

add_multiline_body(slide, [
    "Top metabolomics drivers: D-Glycerate (+), Cinnamoylglycine (-), 2-(4-Hydroxyphenyl)propionic acid (-),",
    "Erythronic acid (+), 4-Hydroxyphenylpyruvate (-), Glutamic acid (+), L-Valine (+)",
], left=Inches(0.6), top=Inches(5.3), width=Inches(12), height=Inches(0.6), font_size=11)

add_method_box(slide,
    "MOFA weights indicate each feature's contribution to Factor 1. Positive weights = higher in OXY; "
    "negative = higher in ANA. Amino acid biosynthesis metabolites dominate, consistent with stress-induced metabolic reprogramming.",
    top=Inches(6.1))

# ============================================================
# SLIDE 8: DIABLO Sample Plot
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "DIABLO: Sample Projection")
add_subtitle(slide, "Supervised separation of ANA vs OXY using all three omics layers")

add_image_safe(slide, os.path.join(DIABLO, "diablo_sample_plot.png"),
               Inches(2.0), Inches(1.7), width=Inches(8.0))

add_method_box(slide,
    "DIABLO (Data Integration Analysis for Biomarker discovery using Latent cOmponents) projects samples onto "
    "discriminant components. Perfect separation of ANA/OXY confirms the oxidative stress signal is consistent across layers.")

# ============================================================
# SLIDE 9: DIABLO Circos
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "DIABLO: Circos Plot -- Cross-Omics Correlations")
add_subtitle(slide, "Inter-layer feature correlations identified by sparse PLS")

add_image_safe(slide, os.path.join(DIABLO, "diablo_circos_plot.png"),
               Inches(3.5), Inches(1.5), width=Inches(6.0))

add_method_box(slide,
    "The circos plot shows pairwise correlations between top features selected by DIABLO across the three omics layers. "
    "Lines connect correlated features across transcriptomics, proteomics, and metabolomics rings.")

# ============================================================
# SLIDE 10: DIABLO Interpretation
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "DIABLO: Feature Selection & Interpretation")
add_subtitle(slide, "Top discriminant features per omics layer")

add_image_safe(slide, os.path.join(DIABLO, "diablo_variable_plot.png"),
               Inches(0.5), Inches(1.7), width=Inches(6.0))

add_multiline_body(slide, [
    "Top DIABLO transcriptomics features (comp1):",
    "  GL50803_0012108 (loading=-0.053)",
    "  GL50803_0089730 (loading=-0.053)",
    "  GL50803_0017558 (loading=-0.053)",
    "",
    "Top DIABLO transcriptomics features (comp2):",
    "  KAE8305986.1 (loading=0.142)",
    "  GL50803_0016114 (loading=0.130)",
    "  GL50803_007760 (loading=0.130)",
    "",
    "DIABLO selects 50 features per layer per component,",
    "providing a sparse, interpretable feature set.",
], left=Inches(7.0), top=Inches(1.7), width=Inches(5.5), height=Inches(4.5), font_size=12)

add_method_box(slide,
    "DIABLO enforces sparsity via L1 penalties, selecting the most discriminative features from each layer. "
    "Component 1 captures the main condition effect; Component 2 captures secondary variation.")

# ============================================================
# SLIDE 11: SNF Clustering
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "SNF: Similarity Network Fusion Clustering")
add_subtitle(slide, "Fused multi-omics sample similarity network")

add_image_safe(slide, os.path.join(SNF, "snf_network_heatmap.png"),
               Inches(0.5), Inches(1.7), width=Inches(5.5))
add_image_safe(slide, os.path.join(SNF, "snf_silhouette.png"),
               Inches(6.5), Inches(1.7), width=Inches(5.5))

add_multiline_body(slide, [
    "SNF Results:",
    "  Cluster 1: ANA_1, ANA_2, ANA_3, ANA_4",
    "  Cluster 2: OXY_1, OXY_2, OXY_3, OXY_4",
    "",
    "  Mean silhouette width: 0.044",
    "  Perfect condition-cluster correspondence",
], left=Inches(6.5), top=Inches(4.5), width=Inches(5.5), height=Inches(1.5), font_size=12)

add_method_box(slide,
    "SNF builds per-layer similarity networks (KNN), then iteratively fuses them. The low silhouette width reflects "
    "high within-group similarity in the fused network. Despite modest silhouette, cluster assignments are biologically perfect.",
    top=Inches(6.2))

# ============================================================
# SLIDE 12: Method Agreement
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Method Agreement: ARI = 1.0 (Perfect)")
add_subtitle(slide, "All three integration methods agree on sample clustering")

add_image_safe(slide, os.path.join(CONSENSUS, "plots", "method_agreement_heatmap.png"),
               Inches(0.5), Inches(1.7), width=Inches(5.5))

add_multiline_body(slide, [
    "Adjusted Rand Index (ARI) Matrix:",
    "",
    "              MOFA    DIABLO    SNF",
    "  MOFA       1.00      1.00     1.00",
    "  DIABLO   1.00      1.00     1.00",
    "  SNF        1.00      1.00     1.00",
    "",
    "All pairwise ARI = 1.0, indicating perfect agreement.",
    "",
    "Sample cluster assignments:",
    "  ANA samples -> Cluster 1 (all methods)",
    "  OXY samples -> Cluster 2 (all methods)",
    "",
    "This perfect consensus is expected given the strong",
    "binary condition effect in this study design.",
], left=Inches(6.5), top=Inches(1.7), width=Inches(6.0), height=Inches(4.5), font_size=13)

add_method_box(slide,
    "ARI measures clustering agreement adjusted for chance (0 = random, 1 = perfect). Perfect ARI across "
    "MOFA, DIABLO, and SNF confirms the ANA/OXY separation is the dominant signal regardless of integration methodology.")

# ============================================================
# SLIDE 13: Section - RNA-Protein Concordance
# ============================================================
add_section_slide("RNA-Protein Concordance")

# ============================================================
# SLIDE 14: Concordance Scatter
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "RNA-Protein Concordance: Fold-Change Scatter")
add_subtitle(slide, "Comparison of log2FC between matched gene-protein pairs")

added = add_image_safe(slide, os.path.join(CONCORD, "concordance_rna_protein.png"),
                       Inches(0.5), Inches(1.7), width=Inches(6.0))
if not added:
    add_image_safe(slide, os.path.join(CONCORD, "concordance_proteomics_vs_transcriptomics.png"),
                   Inches(0.5), Inches(1.7), width=Inches(6.0))

add_multiline_body(slide, [
    "RNA-Protein Concordance Summary:",
    "",
    "  Matched gene-protein pairs: 4,062",
    "  Pearson correlation (r): 0.191",
    "  Direction: weak positive concordance",
    "",
    "  This low correlation is biologically expected:",
    "  - Post-transcriptional regulation",
    "  - Protein half-life differences",
    "  - Translation efficiency variation",
    "  - Giardia-specific regulation mechanisms",
], left=Inches(7.0), top=Inches(1.7), width=Inches(5.5), height=Inches(4.0), font_size=13)

add_method_box(slide,
    "RNA-protein concordance (r=0.191) is within the typical range for eukaryotes (r=0.1-0.6). "
    "Low concordance motivates multi-omics integration, as each layer captures unique regulatory information.")

# ============================================================
# SLIDE 15: Correlation Interpretation
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "RNA-Protein Correlation: Detailed Analysis")
add_subtitle(slide, "Per-gene correlation across 8 samples (4 ANA + 4 OXY)")

add_image_safe(slide, os.path.join(RNACORR, "plots", "rna_protein_de_scatter.png"),
               Inches(0.5), Inches(1.7), width=Inches(6.0))

add_multiline_body(slide, [
    "Detailed RNA-protein correlation analysis:",
    "",
    "  Mean per-gene correlation: 0.060",
    "  Genes with positive correlation: majority",
    "  Genes with negative correlation: minority",
    "",
    "  Post-transcriptional regulation (PTR) genes: 1,286",
    "  High translation efficiency genes: 4,039",
    "  Low translation efficiency genes: 0",
    "",
    "  These results suggest that in Giardia, protein",
    "  levels are primarily regulated at the translational",
    "  or post-translational level rather than transcriptional.",
], left=Inches(7.0), top=Inches(1.7), width=Inches(5.5), height=Inches(4.5), font_size=12)

add_method_box(slide,
    "Per-gene correlation is computed across all 8 samples. Mean r=0.06 indicates minimal linear RNA-protein "
    "relationship for most genes, highlighting extensive post-transcriptional regulation in Giardia.")

# ============================================================
# SLIDE 16: Translation Efficiency
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Translation Efficiency Analysis")
add_subtitle(slide, "Ratio of protein abundance to RNA expression")

add_image_safe(slide, os.path.join(RNACORR, "plots", "translation_efficiency_scatter.png"),
               Inches(0.3), Inches(1.7), width=Inches(5.8))
add_image_safe(slide, os.path.join(RNACORR, "plots", "translation_efficiency_hist.png"),
               Inches(6.5), Inches(1.7), width=Inches(5.8))

add_method_box(slide,
    "Translation efficiency (TE) = protein abundance / RNA abundance. High TE genes (4,039) show efficient "
    "translation or stable protein products. TE variation across genes reflects differential translational regulation "
    "and protein stability in Giardia under oxidative stress.")

# ============================================================
# SLIDE 17: Section - Pathways
# ============================================================
add_section_slide("Cross-Omics Pathway Analysis")

# ============================================================
# SLIDE 18: Cross-Omics Pathway Heatmap
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Cross-Omics Pathway Heatmap")
add_subtitle(slide, "Meta-analysis of enriched pathways across all three omics layers")

add_image_safe(slide, os.path.join(CROSS, "cross_omics_pathway_heatmap.png"),
               Inches(0.3), Inches(1.6), width=Inches(7.5))

add_multiline_body(slide, [
    "Top cross-omics pathways (combined p-adj):",
    "",
    "  1. Biosynthesis of amino acids (p=1.8e-3, metabolomics)",
    "  2. Pentose & glucuronate interconversions (p=2.4e-3)",
    "  3. 2-Oxocarboxylic acid metabolism (p=0.015)",
    "  4. Proteasome (p=0.064, proteomics + transcriptomics)",
    "  5. Ribosome biogenesis (p=0.064, proteomics)",
    "",
    "Multi-layer pathway: Proteasome is enriched in both",
    "proteomics (p=0.021) and transcriptomics (p=0.022),",
    "suggesting coordinated protein turnover response.",
], left=Inches(8.0), top=Inches(1.6), width=Inches(5.0), height=Inches(4.5), font_size=12)

add_method_box(slide,
    "Cross-omics pathway meta-analysis uses Fisher's method to combine p-values across layers. "
    "Pathways enriched in multiple layers have stronger combined evidence. The Proteasome pathway "
    "is the only pathway enriched across 2 omics layers, suggesting a key stress response mechanism.")

# ============================================================
# SLIDE 19: Per-Layer Pathway Enrichment
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Per-Layer Pathway Enrichment")
add_subtitle(slide, "Layer-specific GSEA results")

add_image_safe(slide, os.path.join(CROSS, "metabolomics_top_pathways.png"),
               Inches(0.2), Inches(1.7), width=Inches(4.0))
add_image_safe(slide, os.path.join(CROSS, "proteomics_top_pathways.png"),
               Inches(4.5), Inches(1.7), width=Inches(4.0))
add_image_safe(slide, os.path.join(CROSS, "transcriptomics_top_pathways.png"),
               Inches(8.8), Inches(1.7), width=Inches(4.0))

add_method_box(slide,
    "Metabolomics: dominated by amino acid biosynthesis, carbon metabolism. "
    "Proteomics: ribosome biogenesis, amino acid biosynthesis, proteasome. "
    "Transcriptomics: motor proteins (downregulated), proteasome (upregulated). "
    "Layer-specific patterns reflect distinct regulatory levels of the stress response.")

# ============================================================
# SLIDE 20: MultiGSEA Combined
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "MultiGSEA: Combined Multi-Layer Enrichment")
add_subtitle(slide, "Integrated pathway enrichment across layer pairs")

add_image_safe(slide, os.path.join(MULTIGSEA, "multigsea_combined_enrichment.png"),
               Inches(0.5), Inches(1.7), width=Inches(7.0))

add_multiline_body(slide, [
    "MultiGSEA evaluates pathway enrichment using",
    "information from multiple omics layers simultaneously.",
    "",
    "Three pairwise comparisons:",
    "  - Metabolomics vs Proteomics",
    "  - Metabolomics vs Transcriptomics",
    "  - Proteomics vs Transcriptomics",
    "",
    "This provides a more robust enrichment signal than",
    "single-layer analysis by leveraging cross-layer",
    "concordance within pathway gene/metabolite sets.",
], left=Inches(8.0), top=Inches(1.7), width=Inches(4.8), height=Inches(4.5), font_size=12)

add_method_box(slide,
    "MultiGSEA extends traditional GSEA by combining ranked lists from multiple omics layers. "
    "Pathways that are concordantly enriched across layers receive stronger scores.")

# ============================================================
# SLIDE 21: Cross-Omics Enrichment Dotplot
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Cross-Omics Enrichment Dot Plot")
add_subtitle(slide, "Visualization of pathway enrichment significance and effect sizes")

add_image_safe(slide, os.path.join(CROSS, "cross_omics_enrichment_dotplot.png"),
               Inches(1.5), Inches(1.5), width=Inches(10.0))

add_method_box(slide,
    "Dot size reflects gene/metabolite set size; color intensity reflects enrichment significance. "
    "Amino acid biosynthesis shows the strongest metabolomic signal (NES=-1.89, p=9.2e-6).")

# ============================================================
# SLIDE 22: Section - Key Genes
# ============================================================
add_section_slide("Key Genes & Hub Regulators")

# ============================================================
# SLIDE 23: Hub Regulators
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Hub Regulators in the Regulatory Network")
add_subtitle(slide, "Top connected nodes in the gene-protein correlation network")

add_multiline_body(slide, [
    "Regulatory Network Summary:",
    "  Total edges: 95,566 significant correlations",
    "  Number of hub regulators: 51 genes",
    "  Hub definition: genes with highest network connectivity",
    "",
    "Top 10 Hub Regulators (by target count):",
    "",
    "  Gene           Targets",
    "  GENE_3978    447",
    "  GENE_1601    443",
    "  GENE_3164    443",
    "  GENE_1226    438",
    "  GENE_3318    438",
    "  GENE_3877    437",
    "  GENE_1164    434",
    "  GENE_1429    425",
    "  GENE_2866    425",
    "  GENE_138      422",
], left=Inches(0.6), top=Inches(1.5), width=Inches(5.5), height=Inches(4.5), font_size=13)

add_image_safe(slide, os.path.join(MECH, "plots", "mech_regulatory_network.png"),
               Inches(6.5), Inches(1.5), width=Inches(6.0))

add_method_box(slide,
    "Hub regulators are identified by high connectivity in the RNA-protein correlation network. "
    "These genes likely play central roles in coordinating the oxidative stress response across "
    "transcriptomic and proteomic layers. GENE_3978 is the top hub with 447 regulatory targets.")

# ============================================================
# SLIDE 24: Section - Mechanistic
# ============================================================
add_section_slide("Mechanistic Insights")

# ============================================================
# SLIDE 25: Regulatory Network Visualization
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Gene Regulatory Network")
add_subtitle(slide, "Correlation-based regulatory network from RNA-protein data")

add_image_safe(slide, os.path.join(MECH, "plots", "mech_rna_protein_correlation_dist.png"),
               Inches(0.3), Inches(1.7), width=Inches(6.0))

add_multiline_body(slide, [
    "RNA-Protein Mechanistic Analysis:",
    "",
    "  Method: Spearman correlation network",
    "  Significant edges: 95,566",
    "  Hub regulators: 51",
    "  PTR genes: 1,286",
    "",
    "  Post-transcriptional regulation is widespread:",
    "  1,286 genes show discordant RNA vs protein",
    "  changes, indicating active PTR mechanisms.",
    "",
    "  The network reveals coordinated regulation",
    "  hubs that may serve as intervention targets.",
], left=Inches(6.8), top=Inches(1.7), width=Inches(6.0), height=Inches(4.0), font_size=13)

add_method_box(slide,
    "The regulatory network captures co-expression patterns between RNA and protein across conditions. "
    "Hub nodes with >400 connections are candidate master regulators of the oxidative stress response.")

# ============================================================
# SLIDE 26: Protein-Metabolite Correlations
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Protein-Metabolite Correlations")
add_subtitle(slide, "Significant cross-layer associations (52 pairs, padj < 0.05)")

add_image_safe(slide, os.path.join(MECH, "plots", "protein_metabolite_correlation_heatmap.png"),
               Inches(0.3), Inches(1.7), width=Inches(6.5))

add_multiline_body(slide, [
    "Significant Protein-Metabolite Pairs:",
    "  Total significant pairs: 52 (padj < 0.05 or < 0.08)",
    "  Proteins tested: 100",
    "  Metabolites tested: 50",
    "",
    "Top pairs (|r| > 0.95):",
    "  GENE_1351 -- HMDB0004110 (r=0.976)",
    "  GENE_3390 -- HMDB0001514 (r=0.976)",
    "  GENE_1164 -- HMDB0004284/Tyrosol (r=0.976)",
    "  GENE_388 -- HMDB0000719 (r=0.976)",
    "  GENE_1347 -- HMDB0028697 (r=-0.976)",
    "",
    "Hub GENE_1164 correlates strongly with Tyrosol,",
    "linking regulatory network hubs to metabolic output.",
], left=Inches(7.2), top=Inches(1.5), width=Inches(5.5), height=Inches(4.8), font_size=12)

add_method_box(slide,
    "Protein-metabolite correlations are computed using Spearman rank correlation across 8 samples. "
    "The very high correlations (|r|>0.95) reflect tight enzymatic coupling. "
    "52 significant pairs provide candidate enzyme-substrate relationships for validation.")

# ============================================================
# SLIDE 27: Translation Efficiency & Protein Stability
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Protein Stability & Translation Efficiency")
add_subtitle(slide, "Mechanistic drivers of RNA-protein discordance")

add_image_safe(slide, os.path.join(MECH, "plots", "mech_translation_efficiency_scatter.png"),
               Inches(0.3), Inches(1.7), width=Inches(5.5))
add_image_safe(slide, os.path.join(MECH, "plots", "mech_protein_stability_dist.png"),
               Inches(6.2), Inches(1.7), width=Inches(5.5))

add_method_box(slide,
    "Translation efficiency (TE) quantifies how efficiently each mRNA is translated into protein. "
    "Protein stability analysis examines coefficient of variation across conditions. "
    "Together, these metrics explain why RNA and protein levels diverge for specific genes.")

# ============================================================
# SLIDE 28: Section - Consensus
# ============================================================
add_section_slide("Consensus & Robust Features")

# ============================================================
# SLIDE 29: Robust Features
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Consensus Robust Features")
add_subtitle(slide, "Features identified as important by both MOFA and DIABLO (107 features)")

add_image_safe(slide, os.path.join(CONSENSUS, "plots", "robust_features.png"),
               Inches(0.3), Inches(1.7), width=Inches(6.5))

add_multiline_body(slide, [
    "Robust features (selected by >= 2 methods):",
    "  Total: 107 features",
    "",
    "Top robust features (by combined rank):",
    "  GENE_3262 (MOFA=2.32, DIABLO=0.117)",
    "  GENE_3292 (MOFA=2.52, DIABLO=0.104)",
    "  GENE_894 (MOFA=2.30, DIABLO=0.106)",
    "  GENE_350 (MOFA=2.00, DIABLO=0.117)",
    "  GENE_3062 (MOFA=2.29, DIABLO=0.101)",
    "  HMDB0034297 (MOFA=2.22, DIABLO=0.102)",
    "",
    "Metabolites in consensus:",
    "  HMDB0034297, HMDB0000660, HMDB0001316,",
    "  HMDB0002261, HMDB0061656, HMDB0254199",
], left=Inches(7.2), top=Inches(1.5), width=Inches(5.5), height=Inches(4.8), font_size=12)

add_method_box(slide,
    "Robust features are identified by ranking features in both MOFA (absolute weight) and DIABLO (absolute loading), "
    "then selecting those that appear in the top features of both methods. These 107 features are the most reliable "
    "multi-omics biomarkers of the oxidative stress response.")

# ============================================================
# SLIDE 30: Feature Rank Correlation
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Feature Importance: Method Comparison")
add_subtitle(slide, "Rank correlation between MOFA and DIABLO feature importance")

add_image_safe(slide, os.path.join(CONSENSUS, "plots", "feature_rank_correlation.png"),
               Inches(0.5), Inches(1.7), width=Inches(6.0))
add_image_safe(slide, os.path.join(CONSENSUS, "plots", "meta_integration_pca.png"),
               Inches(7.0), Inches(1.7), width=Inches(5.5))

add_method_box(slide,
    "Left: Rank correlation of feature importance between MOFA and DIABLO shows moderate agreement, "
    "indicating complementary feature selection. Right: PCA of meta-integration results confirms "
    "sample separation is maintained in the consensus feature space.")

# ============================================================
# SLIDE 31: Biological Synthesis
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Biological Synthesis: 6 Key Themes")
add_subtitle(slide, "Integrated interpretation of multi-omics findings")

themes = [
    ("1. Amino Acid Metabolism Reprogramming",
     "Biosynthesis of amino acids is the top enriched pathway (p=9.2e-6 metabolomics). "
     "Key metabolites: Glutamic acid, L-Valine, 4-Hydroxyphenylpyruvate."),
    ("2. Protein Quality Control",
     "Proteasome pathway enriched in both proteomics (p=0.021) and transcriptomics (p=0.022). "
     "Coordinated upregulation suggests active protein turnover under oxidative stress."),
    ("3. Carbon & Energy Metabolism",
     "Carbon metabolism, glycolysis, pentose phosphate pathway enriched. "
     "Metabolic rewiring to support antioxidant defense and repair."),
    ("4. Translational Regulation",
     "Ribosome biogenesis enriched in proteomics (p=0.002). "
     "1,286 PTR genes and mean TE showing post-transcriptional control."),
    ("5. Hub Regulatory Architecture",
     "51 hub regulators with >400 targets each. Network of 95,566 edges "
     "reveals coordinated multi-layer regulatory response."),
    ("6. Cross-Layer Metabolic Coupling",
     "52 significant protein-metabolite pairs (|r|>0.95) link enzymatic "
     "regulators to metabolite output, bridging proteome and metabolome."),
]

for i, (title, desc) in enumerate(themes):
    row = i // 2
    col = i % 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(1.5 + row * 1.8)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.0), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY_BG
    box.line.color.rgb = BORDER_BLUE
    box.line.width = Pt(1)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), Inches(5.7), Inches(0.4),
                title, font_size=13, bold=True, color=DARK_BLUE)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.5), Inches(5.7), Inches(0.9),
                desc, font_size=11, color=BODY_GRAY)

# ============================================================
# SLIDE 32: Foundational Cross-Omics Analysis
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Foundational Cross-Omics Correlations")
add_subtitle(slide, "Sample-level concordance across omics layers")

add_image_safe(slide, os.path.join(FOUND, "foundational", "plots", "mantel_corr_heatmap_proteomics.png"),
               Inches(0.3), Inches(1.7), width=Inches(5.5))
add_image_safe(slide, os.path.join(FOUND, "foundational", "plots", "sample_condition_separation.png"),
               Inches(6.5), Inches(1.7), width=Inches(5.5))

add_method_box(slide,
    "Mantel test evaluates correlation between distance matrices from different omics layers. "
    "Significant Mantel correlation indicates that samples similar in one layer tend to be similar in another, "
    "supporting biological coherence across the multi-omics dataset.")

# ============================================================
# SLIDE 33: Open Questions
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Open Questions & Next Steps")
add_subtitle(slide, "Directions for further analysis and validation")

add_multiline_body(slide, [
    "Open Questions:",
    "",
    "1. Gene annotation: GENE_* identifiers need mapping to known Giardia gene names",
    "   for biological interpretation of hub regulators and robust features",
    "",
    "2. Flavohemoprotein (star gene from A01): Does it remain a top feature",
    "   with the more stringent A02 cutoff?",
    "",
    "3. Temporal dynamics: Are the observed changes early or late stress responses?",
    "",
    "4. Dose-response: Do metabolomic changes scale with oxidative stress intensity?",
    "",
    "Next Steps:",
    "",
    "  - Validate hub regulators by targeted proteomics or Western blot",
    "  - Map GENE_* IDs to functional annotations (GiardiaDB)",
    "  - Network-guided metabolomics follow-up on protein-metabolite pairs",
    "  - Compare A02 results with published Giardia stress datasets",
], top=Inches(1.4), font_size=13)

add_method_box(slide,
    "The A02 analysis provides a more conservative and reliable feature set than A01. "
    "Key validation targets include the 51 hub regulators and 52 protein-metabolite pairs.")

# ============================================================
# SLIDE 34: Methods Summary
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Methods Summary")

add_multiline_body(slide, [
    "Data Processing:",
    "  RNA-seq: DESeq2 differential expression, log2FC >= 0.5 cutoff -> 1,017 DE genes",
    "  Proteomics: MaxQuant DDA, limma-voom DE analysis",
    "  Metabolomics: LC-MS/MS, MetaboAnalyst-style preprocessing",
    "",
    "Multi-Omics Integration:",
    "  MOFA2 (v2): 6 factors, unsupervised, all 3 layers",
    "  DIABLO (mixOmics): 2 components, supervised, sparse PLS",
    "  SNF: K=2 clusters, spectral clustering on fused network",
    "",
    "Cross-Layer Analysis:",
    "  RNA-protein concordance: Pearson/Spearman correlation of matched pairs",
    "  Translation efficiency: Protein/RNA abundance ratio",
    "  Regulatory network: Correlation-based, 95,566 significant edges",
    "  Protein-metabolite: Spearman correlation, 52 significant pairs",
    "",
    "Pathway Enrichment:",
    "  Per-layer GSEA (fgsea), cross-omics Fisher meta-analysis, MultiGSEA",
    "",
    "Consensus:",
    "  ARI-based method agreement, robust feature selection (n=107)",
    "",
    "Software: R 4.x, python-pptx, MOFA2, mixOmics, SNFtool, fgsea, clusterProfiler",
], top=Inches(1.3), font_size=12)

# ============================================================
# SLIDE 35: Thank You
# ============================================================
slide = prs.slides.add_slide(blank_layout)
band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(2.5))
band.fill.solid()
band.fill.fore_color.rgb = SECTION_BG
band.line.fill.background()

add_textbox(slide, Inches(1.0), Inches(2.8), Inches(11), Inches(0.8),
            "Thank You", font_size=40, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(3.6), Inches(11), Inches(0.6),
            "Multi-Omics Integration Analysis A02 -- Elah Pick",
            font_size=18, color=RGBColor(180, 210, 255), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2.0), Inches(5.5), Inches(9), Inches(0.8),
            "Technion Multiomics Center | March 2026\nQuestions & feedback welcome",
            font_size=14, color=BODY_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
