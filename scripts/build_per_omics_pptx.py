#!/usr/bin/env python3
"""
Build Elah Pick Per-Omics Results PowerPoint presentation.
Covers RNA-seq, Proteomics, Metabolomics, Lipidomics with QC, DE, enrichment, interpretation.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──────────────────────────────────────────────────────────────────
OUTPUT = "/home/multiomics-core/outputs/elah_pick_mo/Elah_Per_Omics_Results.pptx"
RNA   = "/home/multiomics-core/outputs/elah_pick_rna/Results_Elah_Pick_RNA_1/rna"
PROT  = "/home/multiomics-core/outputs/elah_pick_proteomics/Results_Elah_Pick_Proteomics_A02/proteomics"
METAB = "/home/multiomics-core/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics"
LIP   = "/home/multiomics-core/outputs/elah_lipidomics/Results_Elah_Lipidomics_A01/lipidomics"

# ── Colors ─────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0, 51, 102)
DARK_GRAY  = RGBColor(51, 51, 51)
ACCENT     = RGBColor(0, 102, 153)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE      = RGBColor(255, 255, 255)
SECTION_BG = RGBColor(0, 51, 102)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Presentation ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

slide_number_counter = [0]

# ── Helper functions ───────────────────────────────────────────────────────

def add_slide():
    """Add a blank slide and return it."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide_number_counter[0] += 1
    add_slide_number(slide, slide_number_counter[0])
    return slide

def add_slide_number(slide, num):
    """Small slide number bottom-right."""
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT

def add_title_text(slide, text, left=0.5, top=0.2, width=12.3, height=0.7, size=28):
    """Add a title textbox."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Calibri"
    return tf

def add_body_text(slide, text, left=0.5, top=1.0, width=12.3, height=5.5, size=14):
    """Add body text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"
    return tf

def add_multiline_text(slide, lines, left=0.5, top=1.0, width=12.3, height=5.5, size=14, bold_first=False):
    """Add multiple lines of text, one paragraph each."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        if bold_first and i == 0:
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
    return tf

def add_method_box(slide, header, body_lines, left=0.5, top=1.0, width=12.3, height=None):
    """Add a styled method/interpretation box with accent header."""
    if height is None:
        height = 0.45 + 0.28 * len(body_lines)
    # Background rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()
    # Header
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.08),
        Inches(width - 0.3), Inches(0.35)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"
    # Body
    txBox2 = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.4),
        Inches(width - 0.3), Inches(height - 0.5)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
    return shape

def safe_image(slide, path, left, top, width=None, height=None):
    """Add an image if it exists, otherwise add a placeholder box."""
    if os.path.isfile(path):
        kwargs = {}
        if width:  kwargs["width"]  = Inches(width)
        if height: kwargs["height"] = Inches(height)
        slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
            Inches(width or 4), Inches(height or 3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = DARK_GRAY
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Image not found]\n{os.path.basename(path)}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

def add_two_images(slide, path1, path2, top=1.2, height=5.0, label1=None, label2=None, label_top=None):
    """Two images side by side."""
    safe_image(slide, path1, left=0.3,  top=top, width=6.2, height=height)
    safe_image(slide, path2, left=6.8, top=top, width=6.2, height=height)
    if label1:
        lt = label_top or (top + height + 0.05)
        add_body_text(slide, label1, left=0.3, top=lt, width=6.2, height=0.3, size=11)
    if label2:
        lt = label_top or (top + height + 0.05)
        add_body_text(slide, label2, left=6.8, top=lt, width=6.2, height=0.3, size=11)

def section_divider(title, subtitle=""):
    """Dark blue section divider slide."""
    slide = add_slide()
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECTION_BG
    shape.line.fill.background()
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(180, 210, 230)
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER
    return slide

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════════
slide = add_slide()
# Accent bar
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.15)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_title_text(slide, "Elah Pick  --  Per-Omics Analysis Results", left=1.0, top=2.0, width=11.3, size=34)
add_body_text(slide, "Giardia lamblia: Oxidative Stress vs Anaerobic",
              left=1.0, top=3.0, width=11.3, size=22)
add_multiline_text(slide, [
    "Layers: RNA-seq | Proteomics | Metabolomics | Lipidomics",
    "",
    "Technion Multiomics Center",
    "March 2026"
], left=1.0, top=4.0, width=11.3, size=16)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 1: RNA-seq (slides 2-12)
# ═══════════════════════════════════════════════════════════════════════════
rna_diag = os.path.join(RNA, "Diagnostic_plots")
rna_qc   = os.path.join(rna_diag, "QC_post")
rna_enr  = os.path.join(RNA, "Enrichment", "plots")

# Slide 2: Section divider
section_divider("Transcriptomics (RNA-seq)", "DESeq2 differential expression  |  5,205 genes")

# Slide 3: Study overview
slide = add_slide()
add_title_text(slide, "RNA-seq: Study Overview")
add_method_box(slide, "Methodology", [
    "Organism: Giardia lamblia  |  Comparison: Oxidative Stress vs Anaerobic",
    "Quantification: 5,205 genes after filtering",
    "Normalization: TMM followed by log-CPM transformation",
    "Differential expression: DESeq2 (negative binomial GLM)",
    "Thresholds: adjusted p-value < 0.05",
    "Enrichment: fGSEA against KEGG and custom gene sets",
], left=0.5, top=1.2, width=12.3, height=2.8)
add_method_box(slide, "Key Results at a Glance", [
    "2,393 DE genes: 1,209 upregulated in oxidative stress, 1,184 downregulated",
    "46% of all measured genes are differentially expressed -- broad transcriptional rewiring",
    "Only 1 significant enrichment pathway: microtubule motor activity (GO, NES=-2.0, padj=0.006)",
    "No KEGG pathway reached statistical significance",
], left=0.5, top=4.3, width=12.3, height=2.2)

# Slide 4: QC boxplots
slide = add_slide()
add_title_text(slide, "RNA-seq QC: Raw vs Normalized Expression")
add_two_images(slide,
    os.path.join(rna_diag, "raw_boxplot.png"),
    os.path.join(rna_diag, "norm_boxplot.png"),
    top=1.1, height=4.6,
    label1="Raw expression distributions", label2="TMM-logCPM normalized")
add_method_box(slide, "AI Interpretation", [
    "Raw distributions show consistent medians across all 8 samples (4 anaerobic, 4 oxidative stress),",
    "indicating comparable sequencing depth. After TMM-logCPM normalization, medians and IQRs are well-aligned,",
    "confirming successful removal of compositional bias. No sample shows signs of degradation or technical failure.",
    "QC PASS: All samples suitable for downstream differential expression analysis.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 5: Sample correlation + distance
slide = add_slide()
add_title_text(slide, "RNA-seq QC: Sample Correlation & Distance")
add_two_images(slide,
    os.path.join(rna_diag, "sample_correlation_heatmap.png"),
    os.path.join(rna_diag, "sample_distance_heatmap.png"),
    top=1.1, height=4.6,
    label1="Pearson correlation heatmap", label2="Euclidean distance heatmap")
add_method_box(slide, "AI Interpretation", [
    "Samples cluster perfectly by condition in both correlation and distance matrices.",
    "Within-group correlations are high (dark blue blocks on diagonal), while between-group distances",
    "are larger, confirming that the oxidative stress vs anaerobic contrast drives the dominant variation.",
    "No outlier samples detected -- all replicates are concordant within their treatment group.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 6: PCA plots
slide = add_slide()
add_title_text(slide, "RNA-seq: Principal Component Analysis")
add_two_images(slide,
    os.path.join(rna_diag, "PCA_PC1.vs.PC2_labeled.png"),
    os.path.join(rna_diag, "PCA_PC1.vs.PC3_labeled.png"),
    top=1.1, height=5.0,
    label1="PC1 vs PC2 (labeled)", label2="PC1 vs PC3 (labeled)")
add_method_box(slide, "PCA Interpretation", [
    "PCA reduces high-dimensional gene expression to principal axes of variation.",
    "Clear separation of oxidative stress vs anaerobic samples on PC1 confirms the primary biological effect.",
    "PC2/PC3 may capture secondary variation (e.g., batch, replicate differences).",
], left=0.3, top=6.3, width=12.7, height=1.0)

# Slide 7: Expression histograms
slide = add_slide()
add_title_text(slide, "RNA-seq: Expression Distribution Histograms")
safe_image(slide, os.path.join(rna_diag, "rna_histograms_summary.png"),
           left=1.5, top=1.1, width=10.3, height=5.0)
add_method_box(slide, "AI Interpretation", [
    "All samples show near-identical unimodal expression distributions after normalization,",
    "confirming uniform library complexity and sequencing depth. The absence of bimodal peaks or",
    "irregular shoulders rules out RNA degradation, ribosomal contamination, or batch artifacts.",
    "Both conditions (anaerobic and oxidative stress) share the same overall expression landscape.",
], left=0.3, top=6.3, width=12.7, height=1.1)

# Slide 8: Volcano + MA
slide = add_slide()
add_title_text(slide, "RNA-seq: Volcano Plot & MA Plot")
add_two_images(slide,
    os.path.join(rna_qc, "volcano_Oxidative stress_vs_Anaerobic.png"),
    os.path.join(rna_qc, "ma_Oxidative stress_vs_Anaerobic.png"),
    top=1.1, height=4.6,
    label1="Volcano: -log10(padj) vs log2FC", label2="MA: log2FC vs mean expression")
add_method_box(slide, "AI Interpretation", [
    "The volcano plot reveals a symmetric DE pattern: 1,209 genes up and 1,184 genes down under oxidative stress.",
    "Several genes reach extreme significance (padj < 1e-100) with log2FC up to 5.1, indicating strong, reproducible changes.",
    "The MA plot confirms no expression-level bias: fold changes are evenly distributed across all expression levels,",
    "validating that the normalization was effective and DE calls are not driven by abundance artifacts.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 9: DE summary
slide = add_slide()
add_title_text(slide, "RNA-seq: Differential Expression Summary")
add_method_box(slide, "DE Results", [
    "Total genes tested: 5,205",
    "Differentially expressed (padj < 0.05): 2,393 (46.0%)",
    "Upregulated in Oxidative Stress: 1,209 genes",
    "Downregulated in Oxidative Stress: 1,184 genes",
    "",
    "The near-equal split between up and down genes suggests broad transcriptional reprogramming",
    "rather than activation or repression of a specific pathway.",
], left=0.5, top=1.2, width=12.3, height=3.0)
add_method_box(slide, "Biological Interpretation", [
    "Giardia responds to oxidative stress with massive transcriptional rewiring.",
    "Nearly half of all expressed genes change significantly, reflecting the organism's",
    "need to remodel its anaerobic metabolism when exposed to oxygen-derived stressors.",
    "This scale of response is consistent with a fundamental shift in metabolic strategy.",
], left=0.5, top=4.5, width=12.3, height=2.0)

# Slide 10: Heatmaps
slide = add_slide()
add_title_text(slide, "RNA-seq: Top DE Gene Heatmaps")
add_two_images(slide,
    os.path.join(rna_diag, "heatmap_top_DE.png"),
    os.path.join(rna_diag, "samples_rna_heatmap.png"),
    top=1.1, height=4.6,
    label1="Top DE genes (Z-score scaled)", label2="Sample-level expression heatmap")
add_method_box(slide, "AI Interpretation", [
    "The top 100 DE genes show a sharp, binary expression pattern: two large clusters perfectly separate",
    "anaerobic (blue) from oxidative stress (red) samples with no intermediate expression states.",
    "This on/off switching pattern suggests Giardia activates a distinct stress-response transcriptional program",
    "rather than gradually modulating expression. The clean dendrogram split confirms high reproducibility across replicates.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 11: Enrichment plots
slide = add_slide()
add_title_text(slide, "RNA-seq: Pathway Enrichment (fGSEA)")
add_two_images(slide,
    os.path.join(rna_enr, "pathway_Oxidative_stress_vs_Anaerobic_KEGG_fgsea.png"),
    os.path.join(rna_enr, "pathway_Oxidative_stress_vs_Anaerobic_custom_fgsea.png"),
    top=1.1, height=4.6,
    label1="KEGG pathways (fGSEA)", label2="Custom gene sets (fGSEA)")
add_method_box(slide, "AI Interpretation", [
    "KEGG fGSEA: Proteasome (NES~-2.1) and motor proteins (NES~-2.0) trend as the most downregulated",
    "pathways, while spliceosome and endocytosis show positive NES, but no KEGG pathway reaches padj < 0.05.",
    "Custom gene sets: Microtubule motor activity is the only statistically significant pathway (padj = 0.006).",
    "The sparse enrichment despite 2,393 DE genes reflects Giardia's compact, poorly annotated genome.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 12: Enrichment interpretation
slide = add_slide()
add_title_text(slide, "RNA-seq: Enrichment Interpretation")
add_method_box(slide, "Pathway Enrichment Summary", [
    "Only 1 pathway reached statistical significance across all tested gene sets:",
    "",
    "  Microtubule motor activity (GO Molecular Function)",
    "  NES = -2.0  |  adjusted p-value = 0.006",
    "  Direction: downregulated in oxidative stress",
    "",
    "No KEGG pathway reached significance (all padj > 0.05).",
], left=0.5, top=1.2, width=12.3, height=2.8)
add_method_box(slide, "Biological Interpretation", [
    "The selective downregulation of microtubule motor activity may reflect reduced",
    "intracellular trafficking and cytoskeletal dynamics under oxidative stress.",
    "Giardia relies on flagellar motility for attachment; reduced motor activity could",
    "indicate a shift from motility to stress defense.",
    "",
    "The lack of enriched KEGG pathways despite 2,393 DE genes suggests the response is",
    "distributed across many processes rather than concentrated in canonical pathways.",
    "This is common in organisms with compact, partially annotated genomes like Giardia.",
], left=0.5, top=4.3, width=12.3, height=2.8)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 2: PROTEOMICS (slides 13-25)
# ═══════════════════════════════════════════════════════════════════════════
prot_diag = os.path.join(PROT, "Diagnostic_plots")
prot_qc   = os.path.join(prot_diag, "QC_post")
prot_enr  = os.path.join(PROT, "Enrichment", "plots")
prot_ppi  = os.path.join(PROT, "ppi_networks", "plots")

# Slide 13
section_divider("Proteomics (LC-MS/MS)", "limma differential expression  |  2,093 proteins")

# Slide 14: overview
slide = add_slide()
add_title_text(slide, "Proteomics: Study Overview")
add_method_box(slide, "Methodology", [
    "Platform: LC-MS/MS (label-free quantification)",
    "Total proteins detected: 4,117  |  After filtering: 2,093",
    "Missing value imputation: MinProb (MNAR assumption)",
    "Differential expression: limma (empirical Bayes moderated t-test)",
    "Thresholds: adjusted p-value < 0.05",
    "PCA outlier detection: 2 outliers (AHA-3A, AHA-4A) flagged but retained",
    "Enrichment: fGSEA (KEGG, custom) + ORA (up/down separately)",
    "Network analysis: PPI from STRING database",
], left=0.5, top=1.2, width=12.3, height=3.2)
add_method_box(slide, "Key Results at a Glance", [
    "216 DE proteins: ~10% of measured proteome",
    "Star hit: Flavohemoprotein -- 24.2-fold upregulated, padj = 6.1e-14",
    "Near-significant KEGG: amino acid biosynthesis (NES=-1.81, padj=0.06), ribosome biogenesis (NES=+1.68, padj=0.06)",
    "PPI hub proteins: GL50803_14842 (degree 34) and GL50803_95898 (degree 34)",
], left=0.5, top=4.7, width=12.3, height=2.0)

# Slide 15: Pre vs post imputation boxplots
slide = add_slide()
add_title_text(slide, "Proteomics QC: Pre- vs Post-Imputation")
add_two_images(slide,
    os.path.join(prot_diag, "pre_imputation_boxplot.png"),
    os.path.join(prot_diag, "imputed_boxplot.png"),
    top=1.1, height=4.6,
    label1="Before imputation (missing values present)", label2="After MinProb imputation")
add_method_box(slide, "AI Interpretation", [
    "Pre-imputation boxplots show comparable medians across all 8 samples, with missingness patterns typical of",
    "label-free proteomics. After MinProb imputation, distributions retain their original shape with a slight left-tail",
    "extension (imputed low-abundance values). The MNAR assumption is appropriate here: missing proteins are likely",
    "below the detection limit rather than randomly absent. Sample-to-sample variability is minimal -- QC PASS.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 16: Histograms
slide = add_slide()
add_title_text(slide, "Proteomics QC: Intensity Distributions")
add_two_images(slide,
    os.path.join(prot_diag, "protein_histograms_summary.png"),
    os.path.join(prot_diag, "imputed_histograms_samples_summary_w0.2_s1.6.png"),
    top=1.1, height=4.6,
    label1="Pre-imputation histograms", label2="Post-imputation histograms (w=0.2, s=1.6)")
add_method_box(slide, "AI Interpretation", [
    "Pre-imputation histograms show Gaussian-like distributions with consistent peaks across samples,",
    "confirming uniform ionization and detection. Post-imputation histograms reveal a small left shoulder",
    "(the imputed values drawn from w=0.2, s=1.6 parameters). This shoulder is well-separated from the main",
    "distribution, indicating imputed values do not contaminate the biological signal range.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 17: Sample correlation + distance
slide = add_slide()
add_title_text(slide, "Proteomics QC: Sample Correlation & Distance")
add_two_images(slide,
    os.path.join(prot_diag, "sample_correlation_heatmap.png"),
    os.path.join(prot_diag, "sample_distance_heatmap.png"),
    top=1.1, height=4.6,
    label1="Pearson correlation heatmap", label2="Euclidean distance heatmap")
add_method_box(slide, "AI Interpretation", [
    "Correlation heatmap shows two clean blocks: OXY samples (r > 0.99 within group) and ANA samples",
    "(r > 0.99 within group), with lower between-group correlation (~0.98). Samples AHA-3A and AHA-4A show",
    "slightly lower within-group correlation but remain within acceptable range. The distance heatmap mirrors",
    "this pattern, confirming that the oxidative stress proteome is globally distinct from the anaerobic proteome.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 18: PCA
slide = add_slide()
add_title_text(slide, "Proteomics: Principal Component Analysis")
add_two_images(slide,
    os.path.join(prot_diag, "PCA_PC1.vs.PC2_labeled.png"),
    os.path.join(prot_diag, "PCA_PC1.vs.PC3_labeled.png"),
    top=1.1, height=5.0,
    label1="PC1 vs PC2 (labeled)", label2="PC1 vs PC3 (labeled)")
add_method_box(slide, "AI Interpretation", [
    "PC1 captures the oxidative stress vs anaerobic separation, explaining the dominant source of variation.",
    "All OXY samples cluster tightly, while ANA samples show slightly more dispersion (AHA-3A, AHA-4A).",
    "PC3 reveals secondary variation within the anaerobic group, possibly reflecting biological heterogeneity",
    "in basal metabolic state. The clear PC1 separation confirms a robust proteomic stress signature.",
], left=0.3, top=6.3, width=12.7, height=1.1)

# Slide 19: UMAP + outlier note
slide = add_slide()
add_title_text(slide, "Proteomics: UMAP & Outlier Assessment")
safe_image(slide, os.path.join(prot_diag, "UMAP.png"),
           left=0.5, top=1.1, width=6.5, height=5.0)
add_method_box(slide, "Outlier Note", [
    "Two potential outliers detected: AHA-3A and AHA-4A",
    "These samples show slight deviation from their group cluster.",
    "Decision: RETAINED in the analysis.",
    "",
    "Rationale: Removing them did not substantially change DE",
    "results, and the sample size is already limited.",
    "Their inclusion provides more statistical power.",
], left=7.5, top=1.1, width=5.3, height=3.0)
add_method_box(slide, "UMAP Interpretation", [
    "UMAP preserves local neighborhood structure better than PCA.",
    "The clear group separation on UMAP confirms that the",
    "oxidative stress proteome is globally distinct from anaerobic.",
], left=7.5, top=4.4, width=5.3, height=1.5)

# Slide 20: Volcano + MA
slide = add_slide()
add_title_text(slide, "Proteomics: Volcano & MA Plots")
add_two_images(slide,
    os.path.join(prot_qc, "volcano_Oxidativestress_vs_Anaerobic_summary.png"),
    os.path.join(prot_qc, "ma_Oxidativestress_vs_Anaerobic_summary.png"),
    top=1.1, height=5.2,
    label1="Volcano plot", label2="MA plot")
add_method_box(slide, "AI Interpretation", [
    "The volcano plot highlights Flavohemoprotein as a dramatic outlier (log2FC ~4.6, padj = 6.1e-14),",
    "towering above other DE proteins. Most DE proteins show moderate fold changes (|log2FC| < 2),",
    "indicating a selective rather than global proteomic shift. The MA plot shows no intensity-dependent",
    "bias, and the Flavohemoprotein hit is robust across all expression levels.",
], left=0.3, top=6.5, width=12.7, height=0.9)

# Slide 21: DE summary + flavohemoprotein
slide = add_slide()
add_title_text(slide, "Proteomics: DE Summary & Key Hits")
add_method_box(slide, "Differential Expression Results", [
    "Total proteins tested: 2,093",
    "Differentially expressed (padj < 0.05): 216 (10.3%)",
    "",
    "Compared to RNA-seq (46% DE), the proteome shows a more selective response.",
    "This is expected: post-transcriptional regulation filters which mRNA changes",
    "translate into protein-level differences.",
], left=0.5, top=1.2, width=12.3, height=2.5)
add_method_box(slide, "Spotlight: Flavohemoprotein (GL50803_10330)", [
    "Fold change: 24.2x upregulated in oxidative stress  |  padj = 6.1e-14",
    "",
    "Flavohemoprotein is a key oxidative stress defense enzyme in Giardia.",
    "It detoxifies nitric oxide (NO) and other reactive nitrogen/oxygen species.",
    "Its massive upregulation confirms a canonical stress response at the protein level.",
    "This is consistent with Giardia shifting from anaerobic metabolism to active",
    "detoxification of reactive oxygen species.",
], left=0.5, top=4.0, width=12.3, height=2.8)

# Slide 22: Heatmaps
slide = add_slide()
add_title_text(slide, "Proteomics: Top DE Protein Heatmaps")
add_two_images(slide,
    os.path.join(prot_qc, "heatmap_top25_Oxidativestress_vs_Anaerobic.png"),
    os.path.join(prot_qc, "heatmap_top50_Oxidativestress_vs_Anaerobic.png"),
    top=1.1, height=4.6,
    label1="Top 25 DE proteins", label2="Top 50 DE proteins")
add_method_box(slide, "AI Interpretation", [
    "Top 25 DE proteins show a striking unidirectional pattern: nearly all are strongly upregulated (red) in",
    "oxidative stress. This asymmetry contrasts with RNA-seq (balanced up/down), suggesting post-transcriptional",
    "selection favors stress-defense proteins. The top 50 heatmap introduces downregulated hits, revealing a core",
    "set of metabolic enzymes suppressed under stress. Perfect sample clustering validates biological reproducibility.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 23: Enrichment plots
slide = add_slide()
add_title_text(slide, "Proteomics: Pathway Enrichment (fGSEA & ORA)")
add_two_images(slide,
    os.path.join(prot_enr, "pathway_Oxidativestress_vs_Anaerobic_KEGG_fgsea.png"),
    os.path.join(prot_enr, "pathway_Oxidativestress_vs_Anaerobic_custom_fgsea.png"),
    top=1.1, height=4.6,
    label1="KEGG fGSEA", label2="Custom gene sets fGSEA")
add_method_box(slide, "AI Interpretation", [
    "KEGG fGSEA reveals ribosome biogenesis (NES~+1.7) and RNA polymerase (NES~+1.5) as top upregulated pathways,",
    "while biosynthesis of amino acids (NES~-1.8) and purine metabolism (NES~-1.7) are downregulated.",
    "None reach strict significance (padj > 0.05), but the coordinated trend is biologically coherent:",
    "Giardia upregulates protein synthesis machinery while suppressing biosynthetic pathways under stress.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 24: PPI networks
slide = add_slide()
add_title_text(slide, "Proteomics: Protein-Protein Interaction Networks")
add_two_images(slide,
    os.path.join(prot_ppi, "ppi_network.png"),
    os.path.join(prot_ppi, "hub_proteins.png"),
    top=1.1, height=5.0,
    label1="PPI network (STRING-based)", label2="Hub proteins by degree centrality")
add_method_box(slide, "AI Interpretation", [
    "The PPI network (529 edges) reveals a densely interconnected stress-response module.",
    "Top hubs GL50803_14842 and GL50803_95898 (degree 34 each, betweenness ~535 and ~132) likely serve as",
    "central coordinators. Additional hubs GL50803_10193 and GL50803_34141 (degree 32) form a secondary tier.",
    "This scale-free topology suggests targeting hub proteins could maximally disrupt the oxidative stress response.",
], left=0.3, top=6.3, width=12.7, height=1.1)

# Slide 25: PPI interpretation + communities
slide = add_slide()
add_title_text(slide, "Proteomics: Network Communities & Interpretation")
safe_image(slide, os.path.join(prot_ppi, "network_communities.png"),
           left=0.3, top=1.1, width=6.5, height=5.0)
add_method_box(slide, "PPI Interpretation", [
    "Community detection reveals functionally related protein modules.",
    "",
    "Near-significant KEGG enrichments suggest coordinated changes in:",
    "  - Amino acid biosynthesis (NES=-1.81, padj=0.06): downregulated",
    "  - Ribosome biogenesis (NES=+1.68, padj=0.06): upregulated",
    "",
    "The combined signal: under oxidative stress, Giardia may increase",
    "ribosome production while reducing de novo amino acid synthesis,",
    "possibly redirecting resources toward stress-response proteins.",
], left=7.0, top=1.1, width=5.8, height=4.0)
add_method_box(slide, "Network Topology", [
    "The presence of highly-connected hub proteins (degree 34) indicates",
    "a scale-free network architecture typical of biological systems.",
    "Targeting these hubs could maximally disrupt the stress response.",
], left=7.0, top=5.3, width=5.8, height=1.5)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 3: METABOLOMICS (slides 26-35)
# ═══════════════════════════════════════════════════════════════════════════
met_diag = os.path.join(METAB, "Diagnostic_plots")
met_enr  = os.path.join(METAB, "Enrichment")

# Slide 26
section_divider("Metabolomics (LC-MS/MS)", "limma differential expression  |  662 metabolites")

# Slide 27: overview
slide = add_slide()
add_title_text(slide, "Metabolomics: Study Overview")
add_method_box(slide, "Methodology", [
    "Platform: LC-MS/MS untargeted metabolomics",
    "Metabolites measured: 662 after filtering",
    "Normalization: log2 transformation + median centering",
    "Differential expression: limma (empirical Bayes moderated t-test)",
    "Thresholds: adjusted p-value < 0.05",
    "Feature selection: PLS-DA scores + VIP, Random Forest importance",
    "Enrichment: GSEA, ORA, QEA (quantitative enrichment analysis)",
], left=0.5, top=1.2, width=12.3, height=2.8)
add_method_box(slide, "Key Results at a Glance", [
    "54 DE metabolites (8.2%): 49 upregulated in anaerobic, 5 upregulated in oxidative stress",
    "Strong directional bias: anaerobic condition accumulates metabolites",
    "GSEA: biosynthesis of amino acids (NES=1.80, FDR=0.013), pentose interconversions (NES=1.86, FDR=0.018)",
], left=0.5, top=4.3, width=12.3, height=1.8)

# Slide 28: QC boxplots + density
slide = add_slide()
add_title_text(slide, "Metabolomics QC: Raw vs Normalized Distributions")
add_two_images(slide,
    os.path.join(met_diag, "intensity_boxplot_raw.png"),
    os.path.join(met_diag, "intensity_boxplot_normalized.png"),
    top=1.1, height=4.6,
    label1="Raw intensity boxplots", label2="Normalized intensity boxplots")
add_method_box(slide, "AI Interpretation", [
    "Raw boxplots show consistent medians (~23-24 log2 intensity) across all 8 samples with comparable IQRs.",
    "Oxidative stress samples (OXY) show slightly wider whiskers suggesting greater intensity dispersion.",
    "After log2 transformation and median centering, distributions are well-aligned. Some low-intensity outliers",
    "persist (values near 0-10) but these are sparse and do not affect the main distribution. QC PASS.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 29: PCA
slide = add_slide()
add_title_text(slide, "Metabolomics: Principal Component Analysis")
add_two_images(slide,
    os.path.join(met_diag, "PCA_PC1.vs.PC2_labeled.png"),
    os.path.join(met_diag, "PCA_PC1.vs.PC3_labeled.png"),
    top=1.1, height=5.0,
    label1="PC1 vs PC2 (labeled)", label2="PC1 vs PC3 (labeled)")
add_method_box(slide, "AI Interpretation", [
    "PC1 separates the two conditions, confirming that the metabolomic landscape is fundamentally different",
    "between anaerobic and oxidative stress. The metabolome is the functional readout of upstream transcriptomic",
    "and proteomic changes -- its clear separation validates that gene/protein-level changes translate into",
    "measurable metabolic consequences. PC3 captures within-group biological variability.",
], left=0.3, top=6.3, width=12.7, height=1.1)

# Slide 30: Volcano + MA
slide = add_slide()
add_title_text(slide, "Metabolomics: Volcano & MA Plots")
add_two_images(slide,
    os.path.join(met_diag, "volcano_Anaerobic_vs_Oxidative_stress.png"),
    os.path.join(met_diag, "MA_Anaerobic_vs_Oxidative_stress.png"),
    top=1.1, height=5.2,
    label1="Volcano plot", label2="MA plot")
add_method_box(slide, "AI Interpretation", [
    "The volcano plot shows a dramatic directional asymmetry: the vast majority of significant metabolites",
    "are higher in anaerobic conditions (depleted under oxidative stress). This one-sided pattern suggests",
    "oxidative stress drives active metabolite consumption -- likely channeled into ROS detoxification and",
    "antioxidant biosynthesis. The MA plot confirms this is not an artifact of expression level.",
], left=0.3, top=6.5, width=12.7, height=0.9)

# Slide 31: DE summary
slide = add_slide()
add_title_text(slide, "Metabolomics: DE Summary & Interpretation")
add_method_box(slide, "Differential Expression Results", [
    "Total metabolites tested: 662",
    "Differentially expressed (padj < 0.05): 54 (8.2%)",
    "Higher in Anaerobic: 49 metabolites",
    "Higher in Oxidative Stress: 5 metabolites",
], left=0.5, top=1.2, width=12.3, height=2.0)
add_method_box(slide, "Biological Interpretation", [
    "The strong directional bias (91% of DE metabolites higher in anaerobic) suggests that",
    "oxidative stress causes metabolite depletion -- likely through increased consumption by",
    "detoxification pathways and redirected metabolic flux toward antioxidant production.",
    "",
    "This is consistent with the proteomic data showing upregulated Flavohemoprotein and",
    "potentially altered amino acid metabolism.",
], left=0.5, top=3.5, width=12.3, height=2.5)

# Slide 32: Heatmap top DE
slide = add_slide()
add_title_text(slide, "Metabolomics: Top DE Metabolite Heatmap")
safe_image(slide, os.path.join(met_diag, "heatmap_top50_DE.png"),
           left=0.3, top=1.1, width=7.5, height=6.0)
add_method_box(slide, "AI Interpretation", [
    "The top 50 DE metabolites cluster into two sharp blocks:",
    "red (high in anaerobic) and blue (depleted in oxidative",
    "stress). This confirms the metabolite depletion signature.",
    "",
    "Amino acids (tryptophan, phenylalanine), nucleotides,",
    "and biosynthetic intermediates dominate the top hits.",
    "These metabolites are likely consumed by antioxidant",
    "defense and stress-response pathways.",
    "",
    "The few metabolites higher in oxidative stress may be",
    "direct products of ROS damage or stress-specific",
    "biosynthetic outputs.",
], left=8.0, top=1.1, width=5.0, height=4.5)

# Slide 33: PLS-DA + VIP + RF
slide = add_slide()
add_title_text(slide, "Metabolomics: Multivariate Feature Selection")
safe_image(slide, os.path.join(met_diag, "plsda_scores.png"),
           left=0.2, top=1.1, width=4.2, height=3.8)
safe_image(slide, os.path.join(met_diag, "plsda_vip_scores.png"),
           left=4.6, top=1.1, width=4.2, height=3.8)
safe_image(slide, os.path.join(met_diag, "rf_importance.png"),
           left=9.0, top=1.1, width=4.2, height=3.8)
add_body_text(slide, "PLS-DA scores (left), VIP scores (center), Random Forest importance (right)",
              left=0.2, top=5.0, width=12.9, height=0.3, size=11)
add_method_box(slide, "Feature Selection Methods", [
    "PLS-DA: Supervised method maximizing separation between conditions. VIP > 1 marks important features.",
    "Random Forest: Non-linear classifier ranking features by their contribution to classification accuracy.",
    "Consensus features (high in both VIP and RF importance) are the most robust biomarker candidates.",
], left=0.2, top=5.5, width=12.9, height=1.6)

# Slide 34: Enrichment
slide = add_slide()
add_title_text(slide, "Metabolomics: Pathway Enrichment")
add_two_images(slide,
    os.path.join(met_enr, "enrichment_gsea_nes_barplot.png"),
    os.path.join(met_enr, "enrichment_ora_dotplot.png"),
    top=1.1, height=4.6,
    label1="GSEA NES barplot", label2="ORA dotplot")
add_method_box(slide, "AI Interpretation", [
    "GSEA reveals 3 significant pathways (FDR < 0.05): biosynthesis of amino acids (NES=-1.89), pentose/glucuronate",
    "interconversions (NES=-1.90), and 2-oxocarboxylic acid metabolism (NES=-1.87). All have negative NES,",
    "meaning these metabolites are coordinately depleted under oxidative stress. The ORA dotplot confirms",
    "amino acid and carbon metabolism pathways as the most affected -- concordant with proteomic findings.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 35: Enrichment interpretation
slide = add_slide()
add_title_text(slide, "Metabolomics: Enrichment Interpretation")
add_method_box(slide, "Significant GSEA Pathways", [
    "1. Biosynthesis of amino acids  (NES = 1.80, FDR = 0.013)",
    "   Amino acid metabolites are coordinately higher in anaerobic conditions.",
    "   Under oxidative stress, amino acids may be consumed for glutathione synthesis",
    "   and other antioxidant mechanisms.",
    "",
    "2. Pentose and glucuronate interconversions  (NES = 1.86, FDR = 0.018)",
    "   Pentose pathway metabolites accumulate under anaerobic conditions.",
    "   Under oxidative stress, these may be redirected to NADPH production",
    "   (via the pentose phosphate pathway) to fuel antioxidant defense.",
], left=0.5, top=1.2, width=12.3, height=3.5)
add_method_box(slide, "Cross-Layer Concordance", [
    "Proteomics showed near-significant downregulation of amino acid biosynthesis enzymes (NES=-1.81).",
    "Metabolomics confirms the functional consequence: amino acid metabolites are depleted under",
    "oxidative stress. This enzyme-metabolite concordance strengthens the biological interpretation.",
], left=0.5, top=5.0, width=12.3, height=1.6)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 4: LIPIDOMICS (slides 36-48+)
# ═══════════════════════════════════════════════════════════════════════════
lip_diag = os.path.join(LIP, "Diagnostic_plots")

# Slide 36
section_divider("Lipidomics (LC-MS/MS)", "limma differential expression  |  575 lipid species")

# Slide 37: overview
slide = add_slide()
add_title_text(slide, "Lipidomics: Study Overview")
add_method_box(slide, "Methodology", [
    "Platform: LC-MS/MS untargeted lipidomics",
    "Lipid species measured: 575 after filtering",
    "Normalization: log2 transformation + median centering",
    "Differential expression: limma (empirical Bayes moderated t-test)",
    "Thresholds: adjusted p-value < 0.05",
    "Lipid-specific analyses: class composition, chain length/saturation, pathway activity,",
    "    enzyme network, PLS-DA, Random Forest, ROC curves",
], left=0.5, top=1.2, width=12.3, height=2.8)
add_method_box(slide, "Key Results at a Glance", [
    "165 DE lipids (28.7%): 148 higher in oxidative stress, 17 higher in anaerobic",
    "Opposite direction to metabolomics: lipids ACCUMULATE under oxidative stress",
    "PC class under-represented in DE (adj.P=0.01) -- phosphatidylcholines are spared",
    "Pathway: PLD deactivated, SMase deactivated, SMS activated (membrane remodeling)",
], left=0.5, top=4.3, width=12.3, height=2.0)

# Slide 38: QC boxplots + density
slide = add_slide()
add_title_text(slide, "Lipidomics QC: Raw vs Normalized Distributions")
add_two_images(slide,
    os.path.join(lip_diag, "intensity_boxplot_raw.png"),
    os.path.join(lip_diag, "intensity_boxplot_normalized.png"),
    top=1.1, height=4.6,
    label1="Raw intensity boxplots", label2="Normalized intensity boxplots")
add_method_box(slide, "AI Interpretation", [
    "Raw lipid intensities show comparable medians (~24-26 log2) across all 8 samples (4 OXY, 4 ANA).",
    "Oxidative stress samples display slightly higher medians, which is corrected after normalization.",
    "Low-intensity outliers (< 10 log2) are present but sparse. After log2 + median centering, all samples",
    "show uniform distributions, confirming successful normalization. No sample failures detected -- QC PASS.",
], left=0.3, top=5.9, width=12.7, height=1.5)

# Slide 39: PCA + loadings
slide = add_slide()
add_title_text(slide, "Lipidomics: PCA Scores & Loadings")
add_two_images(slide,
    os.path.join(lip_diag, "PCA_PC1.vs.PC2_labeled.png"),
    os.path.join(lip_diag, "pca_loading_plot.png"),
    top=1.1, height=5.0,
    label1="PCA scores (PC1 vs PC2, labeled)", label2="PCA loadings (lipid contributions)")
add_method_box(slide, "AI Interpretation", [
    "PCA scores show complete separation of OXY and ANA groups along PC1. The loadings plot identifies which",
    "lipid species drive this separation -- species far from origin are the strongest discriminators.",
    "This tight clustering indicates the lipidome undergoes a global, reproducible shift under oxidative stress,",
    "consistent with active membrane remodeling rather than random lipid perturbation.",
], left=0.3, top=6.3, width=12.7, height=1.1)

# Slide 40: Volcano + MA
slide = add_slide()
add_title_text(slide, "Lipidomics: Volcano & MA Plots")
add_two_images(slide,
    os.path.join(lip_diag, "volcano_Oxy_vs_Ana.png"),
    os.path.join(lip_diag, "MA_Oxy_vs_Ana.png"),
    top=1.1, height=5.2,
    label1="Volcano plot", label2="MA plot")
add_method_box(slide, "AI Interpretation", [
    "The volcano plot reveals a striking directional asymmetry opposite to metabolomics: most DE lipids are",
    "HIGHER in oxidative stress. This reciprocal pattern (metabolites depleted, lipids accumulated) suggests",
    "metabolic flux is redirected toward lipid biosynthesis under stress. The MA plot confirms no abundance bias.",
    "This lipid accumulation likely reflects active membrane reinforcement against oxidative damage.",
], left=0.3, top=6.5, width=12.7, height=0.9)

# Slide 41: DE summary
slide = add_slide()
add_title_text(slide, "Lipidomics: DE Summary & Direction Analysis")
add_method_box(slide, "Differential Expression Results", [
    "Total lipid species tested: 575",
    "Differentially expressed (padj < 0.05): 165 (28.7%)",
    "Higher in Oxidative Stress: 148 lipids (89.7% of DE)",
    "Higher in Anaerobic: 17 lipids (10.3% of DE)",
], left=0.5, top=1.2, width=12.3, height=2.0)
add_method_box(slide, "Contrasting with Metabolomics", [
    "Metabolomics: 91% of DE metabolites are LOWER in oxidative stress (depleted)",
    "Lipidomics: 90% of DE lipids are HIGHER in oxidative stress (accumulated)",
    "",
    "This reciprocal pattern suggests oxidative stress triggers active membrane remodeling:",
    "  - Metabolic intermediates are consumed to fuel lipid biosynthesis",
    "  - Lipids accumulate as the cell reinforces membrane integrity",
    "  - Membrane lipid composition changes may serve as a physical barrier against oxidative damage",
], left=0.5, top=3.5, width=12.3, height=3.0)

# Slide 42: Heatmap
slide = add_slide()
add_title_text(slide, "Lipidomics: Top DE Lipid Heatmap")
safe_image(slide, os.path.join(lip_diag, "heatmap_top50_DE.png"),
           left=0.3, top=1.1, width=7.5, height=6.0)
add_method_box(slide, "AI Interpretation", [
    "Top 50 DE lipids form two dominant clusters:",
    "",
    "Upper block (blue in ANA, red in OXY):",
    "Phosphatidylethanolamines, phosphatidylglycerols,",
    "ceramides, and sphingomyelins accumulate under",
    "oxidative stress -- consistent with membrane",
    "reinforcement and sphingolipid pathway activation.",
    "",
    "Lower block (red in ANA, blue in OXY):",
    "A small subset of lipids (mainly gangliosides",
    "and specific PIs) decrease, possibly consumed",
    "in signaling cascades triggered by stress.",
], left=8.0, top=1.1, width=5.0, height=4.8)

# Slide 43: Lipid class composition
slide = add_slide()
add_title_text(slide, "Lipidomics: Lipid Class Composition")
safe_image(slide, os.path.join(lip_diag, "lipid_category_donut_Ana.png"),
           left=0.1, top=1.1, width=3.2, height=4.5)
safe_image(slide, os.path.join(lip_diag, "lipid_category_donut_Oxy.png"),
           left=3.4, top=1.1, width=3.2, height=4.5)
safe_image(slide, os.path.join(lip_diag, "lipid_class_absolute.png"),
           left=6.8, top=1.1, width=3.2, height=4.5)
safe_image(slide, os.path.join(lip_diag, "lipid_class_normalized.png"),
           left=10.1, top=1.1, width=3.2, height=4.5)
add_body_text(slide, "Category donut: Anaerobic (left) | Oxidative (center-left) | Absolute class counts (center-right) | Normalized proportions (right)",
              left=0.2, top=5.7, width=12.9, height=0.3, size=11)
add_method_box(slide, "AI Interpretation", [
    "Comparing donut charts between conditions reveals subtle compositional shifts: sphingolipid proportions",
    "increase in oxidative stress (consistent with SMS activation), while glycerolipid distribution changes.",
    "The lipidome is dominated by glycerophospholipids (PC, PE, PI) and sphingolipids (SM, Cer). PC remains the",
    "largest class in both conditions, supporting the 'PC sparing' effect identified in the ORA analysis.",
], left=0.2, top=6.1, width=12.9, height=1.3)

# Slide 44: Lipid class ORA + interpretation
slide = add_slide()
add_title_text(slide, "Lipidomics: Lipid Class Enrichment (ORA)")
safe_image(slide, os.path.join(lip_diag, "lipid_class_ora.png"),
           left=0.3, top=1.1, width=6.5, height=5.0)
add_method_box(slide, "ORA Interpretation", [
    "Lipid class ORA tests whether specific lipid classes are over- or under-",
    "represented among DE lipids compared to the full lipidome.",
    "",
    "Key finding: PC (phosphatidylcholine) is significantly UNDER-REPRESENTED",
    "in DE lipids (adj.P = 0.01).",
    "",
    "PCs are the primary structural lipids of cell membranes. Their stability",
    "despite widespread lipid changes suggests Giardia actively protects",
    "membrane integrity by maintaining PC levels during oxidative stress.",
    "This 'PC sparing' effect is a known membrane defense mechanism.",
], left=7.0, top=1.1, width=5.8, height=4.5)

# Slide 45: Chain profiles
slide = add_slide()
add_title_text(slide, "Lipidomics: Chain Length & Saturation Profiles")
add_two_images(slide,
    os.path.join(lip_diag, "chain_length_distribution.png"),
    os.path.join(lip_diag, "chain_saturation.png"),
    top=1.1, height=5.0,
    label1="Acyl chain length distribution", label2="Chain saturation profile")
add_method_box(slide, "Chain Profile Interpretation", [
    "Chain length and saturation determine membrane fluidity and permeability.",
    "Shifts toward longer, more saturated chains reduce membrane fluidity -- a potential defense against oxidative damage.",
], left=0.3, top=6.3, width=12.7, height=1.0)

# Slide 46: Pathway network + barplot
slide = add_slide()
add_title_text(slide, "Lipidomics: Lipid Pathway Activity")
add_two_images(slide,
    os.path.join(lip_diag, "lipid_pathway_network.png"),
    os.path.join(lip_diag, "lipid_pathway_barplot.png"),
    top=1.1, height=5.0,
    label1="Pathway network (enzyme reactions)", label2="Pathway activity barplot")
add_method_box(slide, "AI Interpretation", [
    "The pathway network maps enzymatic conversions between lipid classes. Red/thick edges = activated reactions,",
    "blue = deactivated. Key: PLD (PC->PA, PE->PA) is strongly deactivated (z=-4.07, padj=0.001), protecting",
    "phospholipids. SMS (Cer->SM, z=+3.43, padj=0.003) is activated while SMase (SM->Cer, z=-3.65) is shut down.",
    "The barplot quantifies these enzyme activities, revealing a coherent membrane defense program.",
], left=0.3, top=6.3, width=12.7, height=1.1)

# Slide 47: Pathway interpretation
slide = add_slide()
add_title_text(slide, "Lipidomics: Pathway Interpretation")
add_method_box(slide, "Lipid Pathway Activity Under Oxidative Stress", [
    "Deactivated pathways:",
    "  PLD (Phospholipase D): PC -> PA, PE -> PA   [DEACTIVATED]",
    "  SMase (Sphingomyelinase): SM -> Cer          [DEACTIVATED]",
    "",
    "Activated pathways:",
    "  SMS (Sphingomyelin synthase): Cer -> SM      [ACTIVATED]",
], left=0.5, top=1.2, width=12.3, height=2.5)
add_method_box(slide, "Biological Interpretation", [
    "The pathway activity pattern reveals a coherent membrane defense strategy:",
    "",
    "1. PHOSPHOLIPID PROTECTION: PLD deactivation prevents degradation of PC and PE",
    "   (the two most abundant membrane phospholipids). This preserves membrane structure.",
    "",
    "2. SPHINGOMYELIN ACCUMULATION: SMS activation converts ceramide to sphingomyelin,",
    "   while SMase deactivation prevents the reverse reaction. Net effect: SM increases.",
    "   SM-rich membranes are more resistant to oxidative damage and have lower permeability.",
    "",
    "3. CERAMIDE REDUCTION: Cer is consumed by SMS and not regenerated by SMase.",
    "   Since ceramides are pro-apoptotic signaling lipids, their reduction may help",
    "   Giardia avoid programmed cell death under stress.",
], left=0.5, top=4.0, width=12.3, height=3.2)

# Slide 48: PLS-DA + VIP + RF + ROC
slide = add_slide()
add_title_text(slide, "Lipidomics: Multivariate Classification & Biomarkers")
safe_image(slide, os.path.join(lip_diag, "plsda_scores.png"),
           left=0.1, top=1.1, width=3.2, height=3.0)
safe_image(slide, os.path.join(lip_diag, "plsda_vip_scores.png"),
           left=3.5, top=1.1, width=3.2, height=3.0)
safe_image(slide, os.path.join(lip_diag, "rf_importance.png"),
           left=6.9, top=1.1, width=3.2, height=3.0)
safe_image(slide, os.path.join(lip_diag, "roc_curves_top_lipids.png"),
           left=10.3, top=1.1, width=3.0, height=3.0)
add_body_text(slide, "PLS-DA scores | VIP scores | Random Forest importance | ROC curves (top discriminating lipids)",
              left=0.1, top=4.2, width=13.1, height=0.3, size=11)
add_method_box(slide, "Biomarker Potential", [
    "Multiple supervised methods converge on the same top-discriminating lipids.",
    "ROC curves indicate strong classification performance, suggesting these lipids",
    "could serve as biomarkers for oxidative stress status in Giardia.",
], left=0.1, top=4.6, width=13.1, height=1.3)

# Slide 49: Lipid-lipid class correlation
slide = add_slide()
add_title_text(slide, "Lipidomics: Lipid Class Correlation & Biomarker Table")
safe_image(slide, os.path.join(lip_diag, "lipid_class_correlation.png"),
           left=0.3, top=1.1, width=6.3, height=4.8)
safe_image(slide, os.path.join(lip_diag, "biomarker_table.png"),
           left=6.8, top=1.1, width=6.3, height=4.8)
add_method_box(slide, "AI Interpretation", [
    "Lipid class correlation heatmap reveals coordinated regulation: sphingolipids (SM, Cer, HexCer) cluster together,",
    "as do glycerophospholipids (PC, PE, PI, PS). Anti-correlations between TG and some phospholipids suggest",
    "metabolic tradeoffs. Top 20 biomarker candidates all achieve AUC = 1.0 and 100% power, with Ganglioside GD3,",
    "PI(18:0/18:0), and TG(18:2/18:1) as the strongest discriminators (p < 1e-6, |log2FC| ~1.86).",
], left=0.3, top=6.1, width=12.7, height=1.3)

# ═══════════════════════════════════════════════════════════════════════════
# CLOSING: Summary slide
# ═══════════════════════════════════════════════════════════════════════════
slide = add_slide()
# Accent bar top
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.15)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_title_text(slide, "Summary: Cross-Omics Overview", top=0.3, size=30)

# Four-column summary
cols = [
    ("RNA-seq", [
        "5,205 genes",
        "2,393 DE (46%)",
        "1,209 up / 1,184 down",
        "1 significant GO pathway",
        "Microtubule motor activity down",
    ]),
    ("Proteomics", [
        "2,093 proteins",
        "216 DE (10.3%)",
        "Flavohemoprotein: 24.2x up",
        "PPI hubs: degree 34",
        "Amino acid biosynthesis down",
    ]),
    ("Metabolomics", [
        "662 metabolites",
        "54 DE (8.2%)",
        "91% higher in anaerobic",
        "Amino acid biosynthesis sig.",
        "Metabolite depletion in OXY",
    ]),
    ("Lipidomics", [
        "575 lipid species",
        "165 DE (28.7%)",
        "90% higher in OXY stress",
        "PC class spared (protected)",
        "SMS up, PLD/SMase down",
    ]),
]

for i, (title, items) in enumerate(cols):
    x = 0.3 + i * 3.25
    add_method_box(slide, title, items, left=x, top=1.2, width=3.05, height=2.8)

add_method_box(slide, "Cross-Layer Themes", [
    "1. Massive transcriptional rewiring (46% DE genes) is filtered to selective protein changes (10% DE proteins).",
    "2. Metabolites are depleted under oxidative stress, while lipids accumulate -- reciprocal flux redistribution.",
    "3. Amino acid metabolism is coordinately downregulated at both enzyme and metabolite levels.",
    "4. Flavohemoprotein (24.2x) anchors the oxidative stress defense at the protein level.",
    "5. Membrane remodeling: sphingomyelin synthesis activated, phospholipid degradation suppressed, PCs protected.",
    "6. Giardia's oxidative stress response is a coordinated, multi-layer survival strategy combining",
    "   metabolic reprogramming, antioxidant protein production, and physical membrane reinforcement.",
], left=0.3, top=4.3, width=12.7, height=2.8)

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {slide_number_counter[0]}")

import os
size_mb = os.path.getsize(OUTPUT) / (1024 * 1024)
print(f"File size: {size_mb:.1f} MB")
