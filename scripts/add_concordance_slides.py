#!/usr/bin/env python3
"""Append RNA-Protein Concordance Visualization slides to existing PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

PPTX_PATH = "/home/multiomics-core/outputs/elah_pick_mo/Elah_MultiOmics_A02_Complete.pptx"
IMG_DIR = "/home/multiomics-core/outputs/elah_pick_mo/Results_Elah_Pick_MO_A02/multiomics/concordance_viz"

# Design constants
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_GRAY_BG = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
FONT_NAME = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation(PPTX_PATH)
blank_layout = prs.slide_layouts[6]  # blank layout


def add_textbox(slide, left, top, width, height, text, font_size=11,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Add a textbox with given properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_method_box(slide, text, left=None, top=None, width=None, height=None):
    """Add a light gray method/interpretation box."""
    if left is None:
        left = Inches(0.5)
    if width is None:
        width = Inches(12.3)
    if top is None:
        top = Inches(5.8)
    if height is None:
        height = Inches(1.5)

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY_BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    # Title line
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "METHOD / INTERPRETATION"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    run.font.name = FONT_NAME

    # Body text
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = text
    run2.font.size = Pt(9)
    run2.font.color.rgb = BLACK
    run2.font.name = FONT_NAME

    return shape


def add_divider_slide(title_text, subtitle_text=""):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(blank_layout)

    # Dark blue background rectangle
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                title_text, font_size=36, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    if subtitle_text:
        add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(1),
                    subtitle_text, font_size=16, bold=False, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

    return slide


def add_plot_slide(title, img_filename, method_text,
                   img_width=None, img_height=None,
                   method_top=None, method_height=None):
    """Add a slide with a title, plot image, and method box."""
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
                title, font_size=20, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.LEFT)

    # Image - calculate dimensions
    img_path = os.path.join(IMG_DIR, img_filename)

    if img_width is None:
        img_width = Inches(12)
    if img_height is None:
        img_height = Inches(5.1)
    if method_top is None:
        method_top = Inches(5.7)
    if method_height is None:
        method_height = Inches(1.6)

    img_left = (SLIDE_W - img_width) // 2
    img_top = Inches(0.65)

    slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)

    # Method box
    add_method_box(slide, method_text, top=method_top, height=method_height)

    return slide


# ============================================================
# Build slides
# ============================================================

# Data stats for method boxes
import csv
csv_path = "/home/multiomics-core/outputs/elah_pick_mo/Results_Elah_Pick_MO_A02/multiomics/concordance/concordance_transcriptomics_vs_proteomics.csv"
with open(csv_path) as f:
    reader = csv.DictReader(f)
    rows = list(reader)

n_total = len(rows)
n_sig_prot = sum(1 for r in rows if r["sig_prot"] == "TRUE")
n_sig_rna = sum(1 for r in rows if r["sig_rna"] == "TRUE")
n_concordant = sum(1 for r in rows if r["concordance"] == "concordant")
n_discordant = sum(1 for r in rows if r["concordance"] == "discordant")
n_rna_only = sum(1 for r in rows if r["concordance"] == "rna_only")
n_prot_only = sum(1 for r in rows if r["concordance"] == "protein_only")
n_not_sig = sum(1 for r in rows if r["concordance"] == "not_sig")
n_sig_up = sum(1 for r in rows if r["sig_prot"] == "TRUE" and float(r["logFC_prot"]) > 0)
n_sig_down = sum(1 for r in rows if r["sig_prot"] == "TRUE" and float(r["logFC_prot"]) < 0)
n_both = sum(1 for r in rows if r["sig_prot"] == "TRUE" and r["sig_rna"] == "TRUE")

lfc_prot = [float(r["logFC_prot"]) for r in rows if r["logFC_prot"]]
lfc_rna = [float(r["logFC_rna"]) for r in rows if r["logFC_rna"]]
n = len(lfc_prot)
mean_p = sum(lfc_prot)/n
mean_r = sum(lfc_rna)/n
cov = sum((p-mean_p)*(r-mean_r) for p,r in zip(lfc_prot, lfc_rna)) / n
sd_p = (sum((p-mean_p)**2 for p in lfc_prot)/n)**0.5
sd_r = (sum((r-mean_r)**2 for r in lfc_rna)/n)**0.5
pearson_r = cov / (sd_p * sd_r) if sd_p*sd_r > 0 else 0

print(f"Stats: total={n_total}, sig_prot={n_sig_prot}, sig_rna={n_sig_rna}, "
      f"concordant={n_concordant}, discordant={n_discordant}, both={n_both}, r={pearson_r:.3f}")

# --- Section divider ---
add_divider_slide(
    "RNA-Protein Concordance Visualization",
    f"Comparing transcriptomic and proteomic changes across {n_total} matched gene-protein pairs"
)

# --- 1. Lollipop All ---
add_plot_slide(
    "Paired Lollipop — All Significant Proteins (Top 100 by |FC|)",
    "lollipop_all.png",
    f"Paired lollipop plots showing protein fold changes (left) and corresponding RNA fold changes (right) for the same genes. "
    f"Genes are sorted by protein FC. If RNA and protein agree (same direction), the gene is concordant. "
    f"Discordance suggests post-transcriptional regulation. "
    f"Of {n_sig_prot} significant proteins (padj < 0.05), {n_concordant} are concordant with RNA, "
    f"{n_discordant} are discordant, and {n_prot_only} show protein-only changes.",
    img_height=Inches(5.0), method_top=Inches(5.7)
)

# --- 2. Lollipop Up ---
add_plot_slide(
    f"Paired Lollipop — Upregulated Proteins (n={n_sig_up})",
    "lollipop_up.png",
    f"Paired lollipop plots showing protein fold changes (left) and corresponding RNA fold changes (right) "
    f"for {n_sig_up} significantly upregulated proteins. Genes are sorted by protein FC magnitude. "
    f"Concordant genes (RNA also up) suggest transcriptional activation drives the protein increase. "
    f"Discordant or protein-only cases suggest post-transcriptional stabilization or increased translation.",
    img_height=Inches(5.0), method_top=Inches(5.7)
)

# --- 3. Lollipop Down ---
add_plot_slide(
    f"Paired Lollipop — Downregulated Proteins (n={n_sig_down})",
    "lollipop_down.png",
    f"Paired lollipop plots showing protein fold changes (left) and corresponding RNA fold changes (right) "
    f"for {n_sig_down} significantly downregulated proteins. Genes sorted by protein FC. "
    f"Concordant genes (RNA also down) suggest reduced transcription drives the protein decrease. "
    f"Discordant or protein-only cases suggest post-transcriptional degradation or reduced translation efficiency.",
    img_height=Inches(5.0), method_top=Inches(5.7)
)

# --- 4. Pie Up ---
add_plot_slide(
    f"RNA Status for Significantly Upregulated Proteins (n={n_sig_up})",
    "pie_up.png",
    f"Proportion of significantly upregulated proteins whose RNA levels also change. "
    f"A large gray (unchanged) fraction indicates post-transcriptional buffering — RNA doesn't follow protein. "
    f"This reveals the extent to which proteome increases are driven by transcription vs post-transcriptional mechanisms. "
    f"Out of {n_sig_up} upregulated proteins, those with matching RNA up are transcriptionally driven.",
    img_width=Inches(7), img_height=Inches(4.8),
    method_top=Inches(5.7)
)

# --- 5. Pie Down ---
add_plot_slide(
    f"RNA Status for Significantly Downregulated Proteins (n={n_sig_down})",
    "pie_down.png",
    f"Proportion of significantly downregulated proteins whose RNA levels also change. "
    f"A large gray (unchanged) fraction indicates post-transcriptional buffering — RNA doesn't follow protein. "
    f"This reveals the extent to which proteome decreases are driven by transcription vs post-transcriptional mechanisms. "
    f"Out of {n_sig_down} downregulated proteins, those with matching RNA down are transcriptionally driven.",
    img_width=Inches(7), img_height=Inches(4.8),
    method_top=Inches(5.7)
)

# --- 6. Stacked bar ---
add_plot_slide(
    "Concordance Category Distribution",
    "stacked_bar.png",
    f"Distribution of all {n_total} gene-protein pairs across concordance categories. "
    f"Concordant ({n_concordant}) = both change same direction. "
    f"Discordant ({n_discordant}) = opposite directions. "
    f"RNA only ({n_rna_only}) = RNA changes but protein does not. "
    f"Protein only ({n_prot_only}) = protein changes but RNA does not. "
    f"Not significant ({n_not_sig}) = no significant change in either layer.",
    img_width=Inches(6), img_height=Inches(5.0),
    method_top=Inches(5.7)
)

# --- 7. Volcano by protein ---
add_plot_slide(
    "Protein Volcano — Colored by Protein Significance",
    "volcano_by_prot.png",
    f"Standard protein volcano plot (significance vs fold change) with points colored by protein status. "
    f"Green = significantly upregulated ({n_sig_up} proteins), pink = significantly downregulated ({n_sig_down} proteins), "
    f"gray = not significant. Dashed line at padj = 0.05. "
    f"Total of {n_sig_prot} proteins pass significance threshold out of {n_total} tested.",
    method_top=Inches(5.7)
)

# --- 8. Volcano by RNA ---
add_plot_slide(
    "Protein Volcano — Colored by RNA Change Status",
    "volcano_by_rna.png",
    f"Standard protein volcano plot with points colored by RNA status. This reveals whether the most "
    f"significant protein changes have corresponding RNA evidence. Green = RNA also up, pink = RNA also down, "
    f"gray = RNA not significantly changed. Points without RNA support (gray in significant protein zone) "
    f"suggest post-transcriptional regulation. {n_both} of {n_sig_prot} significant proteins also have significant RNA changes.",
    method_top=Inches(5.7)
)

# --- 9. Venn ---
add_plot_slide(
    "Overlap: Significant RNA vs Significant Protein",
    "venn.png",
    f"Overlap between significantly differentially expressed genes at the RNA level ({n_sig_rna}) vs significantly "
    f"changed proteins ({n_sig_prot}). {n_both} genes are significant in both layers. "
    f"Small overlap indicates extensive post-transcriptional regulation. "
    f"RNA-only significant: {n_sig_rna - n_both}, Protein-only significant: {n_sig_prot - n_both}.",
    img_width=Inches(8), img_height=Inches(4.8),
    method_top=Inches(5.7)
)

# --- 10. Correlation scatter ---
add_plot_slide(
    "RNA vs Protein Fold-Change Correlation",
    "correlation_scatter.png",
    f"Direct comparison of fold changes between RNA and protein for all {n_total} matched gene-protein pairs. "
    f"Pearson r = {pearson_r:.3f}. Perfect correlation (r=1) would mean all protein changes are transcriptionally driven. "
    f"The observed low r indicates extensive post-transcriptional control. "
    f"Points are colored by concordance category: concordant (green), discordant (red), RNA/protein only (orange/blue).",
    method_top=Inches(5.7)
)

# Save
prs.save(PPTX_PATH)

total_slides = len(prs.slides)
file_size = os.path.getsize(PPTX_PATH) / (1024 * 1024)
print(f"\nPPTX saved: {PPTX_PATH}")
print(f"Total slides: {total_slides} (was 63, added 11)")
print(f"File size: {file_size:.1f} MB")
