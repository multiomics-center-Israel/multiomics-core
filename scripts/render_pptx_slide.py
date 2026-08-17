#!/usr/bin/env python3
"""Render a .pptx slide to PNG by reconstructing it from the file's own geometry.

Why this exists
---------------
This machine has no LibreOffice, so the usual pptx -> pdf -> png route is
unavailable and slide decks cannot be inspected before they are handed over.
Reviewing layout arithmetic instead is not a substitute: a checker over
positions, overlaps and estimated overflow will happily pass a slide whose text
column ends halfway down while its figure runs to the bottom. Only looking at
the slide catches that.

What it is not
--------------
It is not PowerPoint's renderer. Every position, size, colour and string is read
from the .pptx itself, but the typefaces are substituted (DejaVu Serif for
Georgia, DejaVu Sans for Calibri), so glyph widths and therefore line breaks
shift slightly. Treat the output as a faithful layout proof, not a pixel-exact
preview, and say so when showing it to someone.

Usage
-----
    python scripts/render_pptx_slide.py deck.pptx 3 slide03.png
    python scripts/render_pptx_slide.py deck.pptx --all out_dir/

Requires python-pptx and Pillow, both already present in the base environment.
"""
from __future__ import annotations

import argparse
import io
import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

DPI = 150

FONTS = {
    "serif": "/usr/share/fonts/truetype/dejavu/DejaVuSerif.ttf",
    "serif_bold": "/usr/share/fonts/truetype/dejavu/DejaVuSerif-Bold.ttf",
    "sans": "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "sans_bold": "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
}

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# Typefaces that should map to the serif substitute. Everything else is sans.
SERIF_FACES = ("georgia", "cambria", "palatino", "garamond", "times")

PICTURE_SHAPE_TYPE = 13


def px(emu) -> int:
    return int(round(Emu(emu).inches * DPI))


def load_font(name: str | None, bold: bool, size_pt: float) -> ImageFont.FreeTypeFont:
    face = (name or "").lower()
    kind = "serif" if face.startswith(SERIF_FACES) else "sans"
    if bold:
        kind += "_bold"
    return ImageFont.truetype(FONTS[kind], max(6, int(round(size_pt * DPI / 72))))


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def colour_of(color_obj, default=(0, 0, 0)) -> tuple[int, int, int]:
    if color_obj is None:
        return default
    try:
        rgb = color_obj.rgb
    except Exception:
        return default
    if isinstance(rgb, str):
        return hex_to_rgb(rgb)
    try:
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return default


def slide_background(slide) -> tuple[int, int, int]:
    bg = slide._element.find("{%s}bg" % NS["p"])
    if bg is None:
        return (255, 255, 255)
    srgb = bg.find(".//{%s}srgbClr" % NS["a"])
    return hex_to_rgb(srgb.get("val")) if srgb is not None else (255, 255, 255)


def bullet_char(paragraph) -> str | None:
    props = paragraph._p.find("{%s}pPr" % NS["a"])
    if props is None:
        return None
    char = props.find("{%s}buChar" % NS["a"])
    return char.get("char") if char is not None else None


def wrap_text(draw, text: str, font, max_width: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for word in text.split():
        candidate = (current + " " + word).strip()
        if draw.textlength(candidate, font=font) <= max_width or not current:
            current = candidate
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [""]


def paragraph_style(paragraph):
    """First run wins: pptxgenjs writes one style per paragraph in practice."""
    runs = [r for r in paragraph.runs if r.text]
    if not runs:
        return None
    first = runs[0]
    size = first.font.size.pt if first.font.size else 14
    has_colour = first.font.color is not None and first.font.color.type is not None
    return {
        "size": size,
        "font": load_font(first.font.name, bool(first.font.bold), size),
        "colour": colour_of(first.font.color) if has_colour else (0, 0, 0),
        "text": "".join(r.text for r in runs),
    }


def draw_text_frame(draw, shape, x: int, y: int, w: int, h: int) -> None:
    pad = int(0.08 * DPI)
    blocks = []
    for paragraph in shape.text_frame.paragraphs:
        style = paragraph_style(paragraph)
        if style is None:
            blocks.append(None)
            continue
        text = style["text"]
        bullet = bullet_char(paragraph)
        if bullet:
            text = f"{bullet}  {text}"
        style["lines"] = wrap_text(draw, text, style["font"], w - 2 * pad)
        style["align"] = paragraph.alignment
        blocks.append(style)

    def height(block) -> int:
        if block is None:
            return int(0.10 * DPI)
        return len(block["lines"]) * int(block["size"] * 1.32 * DPI / 72) + int(0.07 * DPI)

    total = sum(height(b) for b in blocks)
    anchor = str(shape.text_frame.vertical_anchor or "")
    if "MIDDLE" in anchor:
        cursor = y + max(0, (h - total) // 2)
    elif "BOTTOM" in anchor:
        cursor = y + max(0, h - total)
    else:
        cursor = y + int(0.02 * DPI)

    for block in blocks:
        if block is None:
            cursor += int(0.10 * DPI)
            continue
        line_h = int(block["size"] * 1.32 * DPI / 72)
        centred = block["align"] is not None and "CENTER" in str(block["align"])
        for line in block["lines"]:
            if centred:
                tx = x + (w - draw.textlength(line, font=block["font"])) / 2
            else:
                tx = x + pad
            draw.text((tx, cursor), line, font=block["font"], fill=block["colour"])
            cursor += line_h
        cursor += int(0.07 * DPI)


def render_slide(slide, width: int, height: int) -> Image.Image:
    canvas = Image.new("RGB", (width, height), slide_background(slide))
    draw = ImageDraw.Draw(canvas)

    for shape in slide.shapes:
        x, y, w, h = px(shape.left), px(shape.top), px(shape.width), px(shape.height)

        if shape.shape_type == PICTURE_SHAPE_TYPE:
            picture = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
            canvas.paste(picture.resize((w, h), Image.LANCZOS), (x, y))
            continue

        if shape.has_table:
            # Tables are drawn as a plain grid: cell text, borders, no styling.
            table = shape.table
            n_rows, n_cols = len(table.rows), len(table.columns)
            widths = [px(c.width) for c in table.columns]
            heights = [px(r.height) for r in table.rows]
            oy = y
            for ri in range(n_rows):
                ox = x
                for ci in range(n_cols):
                    cell = table.cell(ri, ci)
                    draw.rectangle([ox, oy, ox + widths[ci], oy + heights[ri]],
                                   outline=(210, 214, 200), width=1)
                    font = load_font("Calibri", ri == 0, 11)
                    for k, line in enumerate(
                        wrap_text(draw, cell.text, font, widths[ci] - 10)
                    ):
                        draw.text((ox + 6, oy + 5 + k * 18), line, font=font, fill=(30, 40, 30))
                    ox += widths[ci]
                oy += heights[ri]
            continue

        if not shape.has_text_frame:
            continue

        try:
            if shape.fill.type == 1:
                draw.rectangle([x, y, x + w, y + h],
                               fill=colour_of(shape.fill.fore_color, (255, 255, 255)))
        except Exception:
            pass
        try:
            if shape.line.fill.type == 1:
                draw.rectangle([x, y, x + w, y + h],
                               outline=colour_of(shape.line.color, (200, 200, 200)), width=2)
        except Exception:
            pass

        draw_text_frame(draw, shape, x, y, w, h)

    return canvas


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__,
                                     formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", help="path to the .pptx")
    parser.add_argument("slide", nargs="?", help="1-based slide number")
    parser.add_argument("out", nargs="?", help="output PNG (or directory with --all)")
    parser.add_argument("--all", action="store_true", help="render every slide into a directory")
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    width, height = px(prs.slide_width), px(prs.slide_height)
    slides = list(prs.slides)

    if args.all:
        out_dir = args.slide or args.out or "slides"
        os.makedirs(out_dir, exist_ok=True)
        for i, slide in enumerate(slides, 1):
            path = os.path.join(out_dir, f"slide-{i:02d}.png")
            render_slide(slide, width, height).save(path, quality=95)
            print(path)
        return 0

    if args.slide is None or args.out is None:
        parser.error("give a slide number and an output path, or use --all")
    index = int(args.slide)
    if not 1 <= index <= len(slides):
        parser.error(f"slide {index} out of range; the deck has {len(slides)}")
    render_slide(slides[index - 1], width, height).save(args.out, quality=95)
    print(f"{args.out} ({width}x{height}) — layout proof, fonts substituted")
    return 0


if __name__ == "__main__":
    sys.exit(main())
