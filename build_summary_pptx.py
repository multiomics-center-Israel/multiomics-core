#!/usr/bin/env python3
"""Build a summary PPTX for supervisors showing multiomics-core development."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD   = RGBColor(0x1B, 0x2A, 0x4A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)  # sky-400
ACCENT2   = RGBColor(0x81, 0x8C, 0xF8)  # indigo-400
GREEN     = RGBColor(0x4A, 0xDE, 0x80)  # green-400
SUBTLE    = RGBColor(0x94, 0xA3, 0xB8)  # slate-400
ORANGE    = RGBColor(0xFB, 0x92, 0x3C)  # orange-400
PINK      = RGBColor(0xF4, 0x72, 0xB6)  # pink-400

W = Inches(13.333)
H = Inches(7.5)

BASE = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, size=16, color=WHITE, bold=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    return p

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {"left": left, "top": top}
        if width:  kwargs["width"]  = width
        if height: kwargs["height"] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    return False

def card_rect(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def footer(slide, text="Multi-Omics Core  |  Dor Branch Summary  |  2026"):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             text, size=9, color=SUBTLE, alignment=PP_ALIGN.LEFT)


# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_bg(sl)
accent_line(sl, Inches(1.5), Inches(2.3), Inches(10.3), ACCENT)
add_text(sl, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
         "Multi-Omics Core", size=44, color=WHITE, bold=True)
add_text(sl, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
         "Development Summary — Dor Branch", size=28, color=ACCENT)
add_text(sl, Inches(1.5), Inches(4.6), Inches(10), Inches(0.6),
         "Automated pipelines for Proteomics, RNA-seq & Metabolomics",
         size=16, color=SUBTLE)
add_text(sl, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
         "128 files changed  |  42,700+ lines added  |  3 complete pipelines",
         size=14, color=ACCENT2)
footer(sl)


# =========================================================================
# SLIDE 2 — WHAT WAS BUILT (overview)
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
         "What Was Built", size=32, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(1.0), Inches(5), ACCENT)

items = [
    ("3 End-to-End Pipelines",   "Proteomics  ·  RNA-seq  ·  Metabolomics",      ACCENT),
    ("Interactive HTML Reports", "QC, DE, Enrichment, PPI, AI commentary",        GREEN),
    ("Auto-Generated PPTX",     "9-13 slides per run with AI bottom-lines",       ORANGE),
    ("Pipeline Summary HTML",   "Animated dark-theme workflow timeline",           PINK),
    ("Config Wizard (GUI)",     "Browser-based YAML builder with file uploads",    ACCENT2),
    ("CLI Automation (run.R)",  "--wizard  |  --new  |  --config  |  --fresh",     ACCENT),
]

for i, (title, desc, col) in enumerate(items):
    row = i // 2
    col_idx = i % 2
    x = Inches(0.7 + col_idx * 6.2)
    y = Inches(1.4 + row * 1.9)
    card_rect(sl, x, y, Inches(5.8), Inches(1.6))
    accent_line(sl, x + Inches(0.2), y + Inches(0.15), Inches(0.8), col)
    add_text(sl, x + Inches(0.3), y + Inches(0.3), Inches(5.2), Inches(0.5),
             title, size=20, color=WHITE, bold=True)
    add_text(sl, x + Inches(0.3), y + Inches(0.85), Inches(5.2), Inches(0.5),
             desc, size=14, color=SUBTLE)

footer(sl)


# =========================================================================
# SLIDE 3 — PIPELINE ARCHITECTURE (flow)
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
         "Pipeline Architecture", size=32, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(1.0), Inches(5), ACCENT)

steps = [
    ("Config", "wizard.html\nor YAML", ACCENT),
    ("Import &\nFilter", "Auto-detect\nformat", GREEN),
    ("Normalize\n& QC", "PCA, boxplots\ncorrelation", ACCENT2),
    ("Differential\nExpression", "DESeq2 / t-test\nANOVA", ORANGE),
    ("Enrichment\n& Pathway", "GSEA, ORA\nKEGG, GO", PINK),
    ("Outputs", "HTML + PPTX\n+ Excel", GREEN),
]

for i, (label, detail, col) in enumerate(steps):
    x = Inches(0.5 + i * 2.1)
    y = Inches(2.0)
    # Arrow between boxes
    if i > 0:
        arrow_x = x - Inches(0.35)
        add_text(sl, arrow_x, y + Inches(0.5), Inches(0.4), Inches(0.4),
                 "→", size=24, color=SUBTLE, alignment=PP_ALIGN.CENTER)

    card_rect(sl, x, y, Inches(1.8), Inches(1.8))
    accent_line(sl, x + Inches(0.1), y + Inches(0.1), Inches(1.6), col)
    add_text(sl, x + Inches(0.1), y + Inches(0.25), Inches(1.6), Inches(0.7),
             label, size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.1), y + Inches(1.0), Inches(1.6), Inches(0.7),
             detail, size=10, color=SUBTLE, alignment=PP_ALIGN.CENTER)

# Bottom: 3 domains
for i, (domain, col) in enumerate([("Proteomics", ACCENT), ("RNA-seq", GREEN), ("Metabolomics", ORANGE)]):
    x = Inches(1.5 + i * 3.8)
    card_rect(sl, x, Inches(4.5), Inches(3.2), Inches(2.3))
    accent_line(sl, x + Inches(0.2), Inches(4.6), Inches(2.8), col)
    add_text(sl, x + Inches(0.2), Inches(4.75), Inches(2.8), Inches(0.4),
             domain, size=18, color=col, bold=True, alignment=PP_ALIGN.CENTER)

prot_items = "• MaxQuant / DIA-NN input\n• QRILC / MinVal imputation\n• VSN normalization\n• PPI networks (STRING)\n• fdrtool correction"
rna_items = "• DESeq2 factory\n• tximport support\n• Batch correction\n• Deconvolution\n• Binary pattern clustering"
met_items = "• glog10 transform\n• HMDB annotation\n• RF + PLS-DA selection\n• 4 enrichment methods\n• Auto KEGG mapping"

for i, items_txt in enumerate([prot_items, rna_items, met_items]):
    x = Inches(1.5 + i * 3.8)
    add_text(sl, x + Inches(0.2), Inches(5.2), Inches(2.8), Inches(1.5),
             items_txt, size=11, color=SUBTLE)

footer(sl)


# =========================================================================
# SLIDE 4 — PROTEOMICS EXAMPLE OUTPUTS
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "Proteomics — Example Outputs", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(4), ACCENT)

imgs = [
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/Diagnostic_plots/PCA_PC1.vs.PC2.png", "PCA"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/Diagnostic_plots/QC_post/volcano_PFAS_12.5ppm_vs_NT_summary.png", "Volcano Plot"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/ppi_networks/plots/ppi_network.png", "PPI Network"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/Enrichment/plots/pathway_PFAS_12_5ppm_vs_NT_KEGG_fgsea.png", "KEGG Enrichment"),
]

for i, (path, label) in enumerate(imgs):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.5 + col_idx * 6.3)
    y = Inches(1.2 + row_idx * 3.1)
    card_rect(sl, x, y, Inches(6.0), Inches(2.9))
    add_image_safe(sl, path, x + Inches(0.1), y + Inches(0.1),
                   width=Inches(5.8), height=Inches(2.4))
    add_text(sl, x + Inches(0.1), y + Inches(2.5), Inches(5.8), Inches(0.3),
             label, size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

footer(sl)


# =========================================================================
# SLIDE 5 — RNA-seq EXAMPLE OUTPUTS
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "RNA-seq — Example Outputs", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(4), GREEN)

imgs = [
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/rna/Diagnostic_plots/PCA_PC1.vs.PC2.png", "PCA"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/rna/Diagnostic_plots/QC_post/volcano_PFAS_12.5ppm_vs_NT.png", "Volcano Plot"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/rna/Clustering/Hierarchical/Hierarchical_DE_heatmap.png", "DE Heatmap"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/rna/Enrichment/plots/pathway_PFAS_12_5ppm_vs_NT_GO_fgsea.png", "GO Enrichment"),
]

for i, (path, label) in enumerate(imgs):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.5 + col_idx * 6.3)
    y = Inches(1.2 + row_idx * 3.1)
    card_rect(sl, x, y, Inches(6.0), Inches(2.9))
    add_image_safe(sl, path, x + Inches(0.1), y + Inches(0.1),
                   width=Inches(5.8), height=Inches(2.4))
    add_text(sl, x + Inches(0.1), y + Inches(2.5), Inches(5.8), Inches(0.3),
             label, size=11, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

footer(sl)


# =========================================================================
# SLIDE 6 — METABOLOMICS EXAMPLE OUTPUTS
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "Metabolomics — Example Outputs", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(4), ORANGE)

imgs = [
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Diagnostic_plots/PCA_PC1.vs.PC2.png", "PCA"),
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Diagnostic_plots/volcano_Anaerobic_vs_Oxidative_stress.png", "Volcano Plot"),
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Diagnostic_plots/rf_importance.png", "Random Forest"),
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Diagnostic_plots/plsda_scores.png", "PLS-DA Scores"),
]

for i, (path, label) in enumerate(imgs):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.5 + col_idx * 6.3)
    y = Inches(1.2 + row_idx * 3.1)
    card_rect(sl, x, y, Inches(6.0), Inches(2.9))
    add_image_safe(sl, path, x + Inches(0.1), y + Inches(0.1),
                   width=Inches(5.8), height=Inches(2.4))
    add_text(sl, x + Inches(0.1), y + Inches(2.5), Inches(5.8), Inches(0.3),
             label, size=11, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

footer(sl)


# =========================================================================
# SLIDE 7 — METABOLOMICS ENRICHMENT
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "Metabolomics — Enrichment & Feature Selection", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(5), ORANGE)

imgs = [
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Enrichment/enrichment_gsea_nes_barplot.png", "GSEA NES"),
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Enrichment/enrichment_ora_dotplot.png", "ORA Dotplot"),
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Diagnostic_plots/plsda_vip_scores.png", "PLS-DA VIP"),
    (f"{BASE}/outputs/elah_pick_metabolomics/Results_elah_pick_metabolomics_A01/metabolomics/Enrichment/enrichment_qea_barplot.png", "QEA Barplot"),
]

for i, (path, label) in enumerate(imgs):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.5 + col_idx * 6.3)
    y = Inches(1.2 + row_idx * 3.1)
    card_rect(sl, x, y, Inches(6.0), Inches(2.9))
    add_image_safe(sl, path, x + Inches(0.1), y + Inches(0.1),
                   width=Inches(5.8), height=Inches(2.4))
    add_text(sl, x + Inches(0.1), y + Inches(2.5), Inches(5.8), Inches(0.3),
             label, size=11, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

footer(sl)


# =========================================================================
# SLIDE 8 — AUTO-GENERATED PPTX FEATURE
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "Auto-Generated PowerPoint Reports", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(5), PINK)

add_text(sl, Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
         "Each pipeline run automatically generates a presentation-ready PPTX",
         size=16, color=SUBTLE)

# Left card - slide types
card_rect(sl, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0))
accent_line(sl, Inches(0.7), Inches(2.1), Inches(2), ACCENT)
tf = add_text(sl, Inches(0.8), Inches(2.3), Inches(5.3), Inches(0.5),
              "Slides Included (per domain)", size=18, color=WHITE, bold=True)
slides_list = [
    "Title Slide (project, analyst, stats)",
    "Study Design (groups, samples, features)",
    "PCA plot (PC1 vs PC2)",
    "Volcano plots (one per contrast)",
    "MA plots (one per contrast)",
    "Top DE features table",
    "DE Heatmaps",
    "Enrichment plots (GSEA/ORA)",
    "PPI Network (proteomics)",
    "Feature Selection (metabolomics)",
]
for item in slides_list:
    add_para(tf, f"   •  {item}", size=13, color=SUBTLE)

# Right card - AI feature
card_rect(sl, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5))
accent_line(sl, Inches(7.0), Inches(2.1), Inches(2), ACCENT2)
add_text(sl, Inches(7.1), Inches(2.3), Inches(5.3), Inches(0.4),
         "AI-Powered Bottom Lines", size=18, color=WHITE, bold=True)
tf2 = add_text(sl, Inches(7.1), Inches(2.8), Inches(5.3), Inches(1.5),
               "Each slide includes an AI-generated interpretation using Claude Code CLI. "
               "Summarizes key findings in 1-3 concise sentences.",
               size=13, color=SUBTLE)
add_para(tf2, "", size=6, color=SUBTLE)
add_para(tf2, 'Example: "PCA shows clear separation between treatment and control groups '
              'along PC1 (42% variance), suggesting a strong treatment effect."',
         size=12, color=ACCENT2, bold=False)

# Right bottom card - output paths
card_rect(sl, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.2))
accent_line(sl, Inches(7.0), Inches(4.9), Inches(2), GREEN)
add_text(sl, Inches(7.1), Inches(5.1), Inches(5.3), Inches(0.4),
         "Generated Outputs", size=18, color=WHITE, bold=True)
tf3 = add_text(sl, Inches(7.1), Inches(5.6), Inches(5.3), Inches(1.2),
               "proteomics_summary.pptx   (2.4 MB)", size=13, color=SUBTLE)
add_para(tf3, "rnaseq_summary.pptx          (1.1 MB)", size=13, color=SUBTLE)
add_para(tf3, "metabolomics_summary.pptx (4.5 MB)", size=13, color=SUBTLE)

footer(sl)


# =========================================================================
# SLIDE 9 — HTML REPORTS
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "Interactive HTML Reports", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(4), GREEN)

add_text(sl, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
         "Comprehensive RMarkdown reports — one per domain, fully self-contained",
         size=16, color=SUBTLE)

sections = [
    ("Report Sections", ACCENT, [
        "Executive Summary — key metrics at a glance",
        "Analysis Parameters — tabbed (filtering, norm, DE)",
        "Quality Control — PCA, correlation, boxplots",
        "Differential Expression — volcano, MA, heatmaps",
        "Pathway Enrichment — GSEA, ORA (KEGG, GO)",
        "PPI Networks — STRING integration (proteomics)",
        "Feature Selection — RF, PLS-DA (metabolomics)",
        "AI Commentary — collapsible interpretation blocks",
    ]),
    ("Interactive Features", GREEN, [
        "Plotly — zoomable, hoverable plots",
        "DT tables — searchable, sortable, exportable",
        "Tabbed sections — organized by method",
        "Collapsible blocks — AI commentary on demand",
        "Responsive — works on any screen size",
    ]),
    ("Pipeline Summary HTML", PINK, [
        "Dark-themed animated timeline",
        "Step-by-step workflow visualization",
        "Live metrics from pipeline run",
        "Wet Lab → Upstream → Downstream flow",
        "Emoji-coded numbered steps",
    ]),
]

for i, (title, col, items) in enumerate(sections):
    x = Inches(0.5 + i * 4.2)
    card_rect(sl, x, Inches(1.9), Inches(3.9), Inches(5.0))
    accent_line(sl, x + Inches(0.15), Inches(2.0), Inches(1.5), col)
    add_text(sl, x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(0.4),
             title, size=17, color=col, bold=True)
    tf = add_text(sl, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(4.0),
                  "", size=12, color=SUBTLE)
    for item in items:
        add_para(tf, f"  •  {item}", size=12, color=SUBTLE, space_before=Pt(6))

footer(sl)


# =========================================================================
# SLIDE 10 — WIZARD & CLI
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "Configuration & Automation", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(5), ACCENT2)

# Left — Wizard
card_rect(sl, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.5))
accent_line(sl, Inches(0.7), Inches(1.5), Inches(2), ACCENT)
add_text(sl, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5),
         "wizard.html — Browser-Based GUI", size=20, color=WHITE, bold=True)
tf = add_text(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.5),
              "No-code config builder for researchers", size=14, color=SUBTLE)
add_para(tf, "", size=6, color=SUBTLE)
add_para(tf, "  •  Select analysis type (RNA / Proteomics / Metabolomics)", size=13, color=SUBTLE, space_before=Pt(8))
add_para(tf, "  •  Upload data files with auto column detection", size=13, color=SUBTLE, space_before=Pt(8))
add_para(tf, "  •  Interactive contrast builder", size=13, color=SUBTLE, space_before=Pt(8))
add_para(tf, "  •  Configure filtering, normalization, DE method", size=13, color=SUBTLE, space_before=Pt(8))
add_para(tf, "  •  Live YAML preview", size=13, color=SUBTLE, space_before=Pt(8))
add_para(tf, "  •  Download config + Run pipeline from wizard", size=13, color=SUBTLE, space_before=Pt(8))
add_para(tf, "  •  Dark / Light mode toggle", size=13, color=SUBTLE, space_before=Pt(8))
add_para(tf, "", size=6, color=SUBTLE)
add_para(tf, "  Opens with:  Rscript run.R --wizard", size=14, color=ACCENT, bold=True, space_before=Pt(16))

# Right — CLI
card_rect(sl, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5))
accent_line(sl, Inches(7.0), Inches(1.5), Inches(2), GREEN)
add_text(sl, Inches(7.1), Inches(1.7), Inches(5.5), Inches(0.5),
         "run.R — CLI Automation", size=20, color=WHITE, bold=True)
tf2 = add_text(sl, Inches(7.1), Inches(2.3), Inches(5.5), Inches(0.5),
               "Command-line interface for power users", size=14, color=SUBTLE)
add_para(tf2, "", size=6, color=SUBTLE)

cmds = [
    ("Rscript run.R --wizard",  "Open GUI wizard in browser"),
    ("Rscript run.R --new",     "Interactive CLI project setup"),
    ("Rscript run.R --config config.yaml", "Run with existing config"),
    ("Rscript run.R",           "Re-run last config (cached)"),
    ("Rscript run.R --fresh",   "Full re-run (no cache)"),
]
for cmd, desc in cmds:
    add_para(tf2, f"  $ {cmd}", size=12, color=ACCENT, bold=True, space_before=Pt(12))
    add_para(tf2, f"     {desc}", size=11, color=SUBTLE, space_before=Pt(2))

add_para(tf2, "", size=8, color=SUBTLE)
add_para(tf2, "Features:", size=14, color=WHITE, bold=True, space_before=Pt(12))
add_para(tf2, "  •  Auto-detect file format (TSV/CSV)", size=12, color=SUBTLE, space_before=Pt(6))
add_para(tf2, "  •  Smart contrast generation from metadata", size=12, color=SUBTLE, space_before=Pt(6))
add_para(tf2, "  •  targets-based caching (skip unchanged steps)", size=12, color=SUBTLE, space_before=Pt(6))

footer(sl)


# =========================================================================
# SLIDE 11 — PPI & ADVANCED (proteomics extras)
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
         "Proteomics — PPI Networks & Correlation", size=28, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(0.85), Inches(5), ACCENT)

imgs = [
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/ppi_networks/plots/network_communities.png", "Network Communities"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/ppi_networks/plots/hub_proteins.png", "Hub Proteins"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/Diagnostic_plots/sample_correlation_heatmap.png", "Sample Correlation"),
    (f"{BASE}/outputs/amir_sapir/Results_Amir_Sapir_A01/proteomics/Clustering/Hierarchical/Hierarchical_DE_heatmap.png", "Hierarchical Clustering"),
]

for i, (path, label) in enumerate(imgs):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.5 + col_idx * 6.3)
    y = Inches(1.2 + row_idx * 3.1)
    card_rect(sl, x, y, Inches(6.0), Inches(2.9))
    add_image_safe(sl, path, x + Inches(0.1), y + Inches(0.1),
                   width=Inches(5.8), height=Inches(2.4))
    add_text(sl, x + Inches(0.1), y + Inches(2.5), Inches(5.8), Inches(0.3),
             label, size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

footer(sl)


# =========================================================================
# SLIDE 12 — SUMMARY / NEXT STEPS
# =========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
         "Summary & Impact", size=32, color=WHITE, bold=True)
accent_line(sl, Inches(0.7), Inches(1.0), Inches(5), ACCENT)

# Stats row
stats = [
    ("3", "Pipelines", ACCENT),
    ("128", "Files Changed", GREEN),
    ("42,700+", "Lines of Code", ORANGE),
    ("11", "HTML Reports", PINK),
    ("3", "Auto PPTX", ACCENT2),
]
for i, (num, label, col) in enumerate(stats):
    x = Inches(0.5 + i * 2.5)
    card_rect(sl, x, Inches(1.5), Inches(2.2), Inches(1.4))
    add_text(sl, x, Inches(1.6), Inches(2.2), Inches(0.7),
             num, size=36, color=col, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, x, Inches(2.3), Inches(2.2), Inches(0.4),
             label, size=13, color=SUBTLE, alignment=PP_ALIGN.CENTER)

# Key achievements
card_rect(sl, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.5))
accent_line(sl, Inches(0.7), Inches(3.4), Inches(3), ACCENT)
tf = add_text(sl, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5),
              "Key Achievements", size=22, color=WHITE, bold=True)

achievements = [
    "End-to-end automation — from raw data to publication-ready outputs in one command",
    "Three complete pipelines — Proteomics, RNA-seq, Metabolomics — unified architecture",
    "No-code entry point — wizard.html lets researchers configure without coding",
    "AI-enhanced outputs — Claude-powered interpretations in reports and presentations",
    "Production tested — successfully run on multiple real datasets (Amir-Sapir, Elah-Pick)",
]
for a in achievements:
    add_para(tf, f"    •   {a}", size=14, color=SUBTLE, space_before=Pt(10))

footer(sl)


# =========================================================================
# SAVE
# =========================================================================
out_path = f"{BASE}/multiomics_core_summary.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
